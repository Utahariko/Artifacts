from __future__ import annotations

import copy
import csv
import io
import json
import math
import os
import re
import uuid
from datetime import date, datetime, timedelta
from html import escape
from pathlib import Path
from typing import Dict, List, Tuple

import streamlit as st

try:
    from openai import OpenAI
except Exception:  # pragma: no cover - optional dependency at runtime
    OpenAI = None


APP_TITLE = "PlanWise"
SRL_PHASES = ["Planning", "Monitoring", "Control", "Reflection"]
STATUS_OPTIONS = ["Pending", "Completed", "Delayed", "Confused", "Overwhelmed"]
LANGUAGE_OPTIONS = ["English", "中文"]
DATA_DIR = Path("data")
TASKS_FILE = DATA_DIR / "tasks.json"
UPLOAD_DIR = DATA_DIR / "uploads"
SUPPORTED_MATERIAL_TYPES = [
    "txt",
    "md",
    "pdf",
    "docx",
    "pptx",
    "xlsx",
    "csv",
    "png",
    "jpg",
    "jpeg",
    "webp",
]
PERSISTED_TASK_KEYS = [
    "study_plan",
    "plan_context",
    "learning_materials",
    "monitoring_status",
    "monitoring_details",
    "monitoring_feedback",
    "adjusted_plan",
    "control_recommendations",
    "control_next_actions",
    "control_context",
    "control_source",
    "reflection_report",
    "reflection_inputs",
    "coach_history",
    "workflow_current_index",
    "workflow_completed_indices",
    "workflow_step_reflections",
]
TRANSLATIONS = {
    "app_caption": {
        "English": "An LLM-supported Self-Regulated Learning planner for study goals, learning challenges, and academic tasks.",
        "中文": "一个由大语言模型支持的自我调节学习规划工具，用于学习目标、学习困难和学术任务。",
    },
    "phase_planning": {"English": "Planning", "中文": "计划"},
    "phase_monitoring": {"English": "Monitoring", "中文": "监控"},
    "phase_control": {"English": "Control", "中文": "调节"},
    "phase_reflection": {"English": "Reflection", "中文": "反思"},
    "phase_ai_coach": {"English": "AI Coach", "中文": "AI 学习教练"},
    "phase_dashboard": {"English": "Dashboard", "中文": "首页"},
    "phase_task_library": {"English": "Task Library", "中文": "任务库"},
    "phase_task_details": {"English": "Task Details", "中文": "任务详情"},
    "phase_today": {"English": "Today Learning", "中文": "今日学习"},
    "phase_reports": {"English": "Reports", "中文": "历史报告"},
    "phase_guide": {"English": "How to Use", "中文": "使用指南"},
    "phase_desc_planning": {
        "English": "Set goals and transform learning demands into a workable study plan.",
        "中文": "设定目标，并把学习需求转化为可执行的学习计划。",
    },
    "phase_desc_monitoring": {
        "English": "Track progress, confusion, delay, and workload signals.",
        "中文": "追踪进度、困惑、延迟和负荷信号。",
    },
    "phase_desc_control": {
        "English": "Adjust strategies when the plan no longer fits the learning situation.",
        "中文": "当计划不再适合当前学习情况时，调整策略。",
    },
    "phase_desc_reflection": {
        "English": "Review outcomes and build metacognitive awareness for the next cycle.",
        "中文": "回顾学习结果，为下一轮学习建立元认知意识。",
    },
    "nav_title": {"English": "PlanWise Navigation", "中文": "PlanWise 导航"},
    "nav_caption": {"English": "SRL planner stages", "中文": "自我调节学习阶段"},
    "language": {"English": "Language", "中文": "语言"},
    "go_to": {"English": "Go to", "中文": "前往"},
    "dashboard_title": {"English": "Learning Dashboard", "中文": "学习仪表盘"},
    "dashboard_intro": {
        "English": "A web-style overview of your current SRL learning workflow, materials, progress, and next action.",
        "中文": "以网页式仪表盘查看当前自我调节学习流程、材料、进度和下一步行动。",
    },
    "today_learning": {"English": "Today Learning", "中文": "今日学习"},
    "task_details": {"English": "Task Details", "中文": "任务详情"},
    "quick_actions": {"English": "Quick Actions", "中文": "快捷操作"},
    "start_planning": {"English": "Start planning", "中文": "开始计划"},
    "continue_monitoring": {"English": "Continue monitoring", "中文": "继续监控"},
    "open_ai_coach": {"English": "Open AI Coach", "中文": "打开 AI 教练"},
    "open_reports": {"English": "Open reports", "中文": "打开报告"},
    "no_active_task_dashboard": {
        "English": "No active task yet. Start by creating a learning plan or uploading materials.",
        "中文": "当前还没有任务。可以先创建学习计划，或上传学习材料。",
    },
    "task_progress": {"English": "Task Progress", "中文": "任务进度"},
    "materials_count": {"English": "Materials", "中文": "材料数"},
    "reflection_count": {"English": "Reflections", "中文": "反思数"},
    "current_workflow_step": {"English": "Current Workflow Step", "中文": "当前流程步骤"},
    "recent_tasks": {"English": "Recent Tasks", "中文": "最近任务"},
    "reports_title": {"English": "Learning Reports", "中文": "学习报告"},
    "reports_intro": {
        "English": "Review completed steps, reflection records, materials, and SRL profile evidence across the current task.",
        "中文": "回顾当前任务中的已完成步骤、反思记录、材料和自我调节学习画像证据。",
    },
    "no_reports_yet": {
        "English": "No completed step reports yet. Finish a Reflection cycle to create the first report.",
        "中文": "还没有已完成的步骤报告。完成一次反思循环后会生成第一份报告。",
    },
    "progress_timeline": {"English": "Progress Timeline", "中文": "进度时间线"},
    "guide_title": {"English": "How to Use PlanWise", "中文": "如何使用 PlanWise"},
    "guide_intro": {
        "English": "A short teaching guide for using PlanWise as an SRL-supported learning planner rather than a generic chatbot.",
        "中文": "这是一个简短教学指南，说明如何把 PlanWise 当作自我调节学习规划工具使用，而不是普通聊天机器人。",
    },
    "guide_quick_start": {"English": "5-Minute Quick Start", "中文": "5 分钟快速上手"},
    "guide_srl_flow": {"English": "SRL Learning Cycle", "中文": "自我调节学习循环"},
    "guide_class_demo": {"English": "Classroom Demo Script", "中文": "课堂展示脚本"},
    "guide_tips": {"English": "Good Use Tips", "中文": "使用建议"},
    "guide_quick_items": {
        "English": [
            "Open Planning and describe a learning goal, challenge, deadline, difficulty, and available study time.",
            "Optionally upload learning materials such as PDF, DOCX, PPTX, spreadsheet, text, or images.",
            "Generate a study plan and check whether the plan items match the deadline and task scope.",
            "Go to Monitoring after completing the current step and record progress, delay, confusion, or overload.",
            "Use Control to adjust strategy, then Reflection to close the step and unlock the next one.",
        ],
        "中文": [
            "进入计划页面，描述学习目标、当前困难、截止日期、难度和每天可用学习时间。",
            "可选择上传学习材料，例如 PDF、DOCX、PPTX、表格、文本或图片。",
            "生成学习计划，并检查计划步骤是否符合截止日期和任务范围。",
            "完成当前步骤后进入监控，记录进展、延迟、困惑或负荷过高。",
            "使用调节页面调整策略，再进入反思页面完成当前步骤并解锁下一步。",
        ],
    },
    "guide_srl_items": {
        "English": [
            "Planning: turn a learning problem into a structured plan with strategies and time ranges.",
            "Monitoring: record evidence, barriers, confidence, and actual time after each learning step.",
            "Control: adapt the plan when monitoring signals show delay, confusion, or overload.",
            "Reflection: explain what worked, what slowed progress, and what should change in the next cycle.",
        ],
        "中文": [
            "计划：把学习问题转化为带策略和时间范围的结构化计划。",
            "监控：每完成一个学习步骤后，记录证据、障碍、信心和实际用时。",
            "调节：当监控信号显示延迟、困惑或负荷过高时，调整计划和策略。",
            "反思：解释哪些策略有效、什么拖慢了进度，以及下一轮要改变什么。",
        ],
    },
    "guide_demo_items": {
        "English": [
            "Demo 1: Create a task such as preparing for a statistics exam or writing a literature review.",
            "Demo 2: Upload one relevant material and show how PlanWise uses the material brief in planning.",
            "Demo 3: Mark the first step as delayed or confused in Monitoring and generate adaptive feedback.",
            "Demo 4: Open Control to show strategy adjustment, then Reflection to complete the step.",
            "Demo 5: Return to Dashboard or Reports to show progress, task memory, and reflection archive.",
        ],
        "中文": [
            "展示 1：创建一个任务，例如准备统计学考试或写文献综述。",
            "展示 2：上传一份相关材料，说明 PlanWise 如何把材料摘要用于计划生成。",
            "展示 3：在监控中把第一步标记为延迟或困惑，并生成适应性反馈。",
            "展示 4：进入调节页面展示策略调整，再进入反思页面完成该步骤。",
            "展示 5：回到首页或历史报告页，展示进度、任务记忆和反思档案。",
        ],
    },
    "guide_tip_items": {
        "English": [
            "Use specific task descriptions. PlanWise works better with concrete goals than vague intentions.",
            "Treat GenAI output as a planning partner. Check, adapt, and revise the plan instead of accepting it blindly.",
            "Record honest monitoring signals. Delay and confusion are useful data, not failure.",
            "Write short reflections. One useful sentence is better than skipping reflection entirely.",
            "Save tasks when switching projects so previous plans and reports can be reused.",
        ],
        "中文": [
            "尽量具体描述任务。PlanWise 更适合处理明确目标，而不是模糊意图。",
            "把 GenAI 输出当作计划伙伴。要检查、调整和修订，而不是盲目接受。",
            "诚实记录监控信号。延迟和困惑是有用数据，不代表失败。",
            "反思可以很短。一句有用的反思也比完全跳过反思更好。",
            "切换项目时保存任务，这样之前的计划和报告可以继续使用。",
        ],
    },
    "current_task": {"English": "Current Task", "中文": "当前任务"},
    "task_memory": {"English": "Task Memory", "中文": "任务记忆"},
    "task_library_title": {"English": "Task Library", "中文": "任务库"},
    "task_library_intro": {
        "English": "Browse, load, rename, and delete saved learning tasks without disturbing the current task until you explicitly load one.",
        "中文": "在这里浏览、读取、重命名和删除历史学习任务。只有点击读取时，才会切换当前任务。",
    },
    "open_task_library": {"English": "Open Task Library", "中文": "打开任务库"},
    "saved_task_count": {"English": "Saved tasks", "中文": "已保存任务"},
    "new_task": {"English": "New task", "中文": "新建任务"},
    "save_task": {"English": "Save current task", "中文": "保存当前任务"},
    "open_saved_task": {"English": "Open saved task", "中文": "打开历史任务"},
    "load_task": {"English": "Load task", "中文": "读取任务"},
    "rename_task": {"English": "Rename task", "中文": "重命名任务"},
    "task_name": {"English": "Task name", "中文": "任务名称"},
    "task_renamed": {"English": "Task renamed.", "中文": "任务已重命名。"},
    "task_loaded": {"English": "Task loaded.", "中文": "任务已读取。"},
    "delete_task": {"English": "Delete task", "中文": "删除任务"},
    "confirm_delete_task": {"English": "Confirm delete selected task", "中文": "确认删除所选任务"},
    "task_deleted": {"English": "Task deleted.", "中文": "任务已删除。"},
    "select_task_first": {"English": "Select a saved task first.", "中文": "请先选择一个历史任务。"},
    "nothing_to_save": {
        "English": "Nothing to save yet. Create a plan or ask the AI Coach first.",
        "中文": "目前还没有可保存的内容。请先生成计划，或先使用 AI 学习教练。",
    },
    "task_preview": {"English": "Task preview", "中文": "任务预览"},
    "created_at": {"English": "Created", "中文": "创建时间"},
    "updated_at": {"English": "Updated", "中文": "更新时间"},
    "unsaved_task": {"English": "Unsaved task", "中文": "未保存任务"},
    "no_saved_tasks": {"English": "No saved tasks yet", "中文": "暂无历史任务"},
    "task_saved": {"English": "Task saved.", "中文": "任务已保存。"},
    "new_task_ready": {"English": "Started a new blank task.", "中文": "已开始一个新的空白任务。"},
    "autosaved": {"English": "Autosaved", "中文": "已自动保存"},
    "deadline": {"English": "Deadline", "中文": "截止日期"},
    "phase_prefix": {"English": "SRL Phase", "中文": "自我调节学习阶段"},
    "goal_setting": {"English": "Goal Setting", "中文": "目标设定"},
    "planning_intro": {
        "English": "Describe any learning challenge, study goal, or academic task, and PlanWise will turn it into a structured, strategy-supported plan.",
        "中文": "描述任何学习困难、学习目标或学术任务，PlanWise 会将其转化为结构化、带策略支持的学习计划。",
    },
    "learning_goal_or_task": {"English": "Learning goal or task", "中文": "学习目标或任务"},
    "learning_goal_placeholder": {
        "English": "Example: Prepare for a statistics exam, plan a research proposal, improve presentation skills, or understand a difficult course topic.",
        "中文": "例如：准备统计学考试、规划研究计划、提升展示能力，或理解一门课里的难点主题。",
    },
    "perceived_difficulty": {"English": "Perceived difficulty level", "中文": "感知难度"},
    "daily_time": {"English": "Daily available study time", "中文": "每天可用学习时间"},
    "current_challenge": {"English": "Current challenge", "中文": "当前学习困难"},
    "challenge_placeholder": {
        "English": "Example: I keep spending time on materials but still cannot explain the key ideas clearly.",
        "中文": "例如：我花了很多时间看资料，但还是不能清楚解释关键概念。",
    },
    "optional_learning_context": {"English": "Optional learning context", "中文": "可选学习背景"},
    "optional_context_caption": {
        "English": "These details are optional. Add any information that would help PlanWise generate a more personalised SRL plan.",
        "中文": "这些信息都是可选的。你可以补充任何有助于 PlanWise 生成个性化自我调节学习计划的信息。",
    },
    "task_type": {"English": "Task type", "中文": "任务类型"},
    "task_type_placeholder": {
        "English": "Example: exam revision, essay, presentation, project, coding practice",
        "中文": "例如：考试复习、论文、展示、项目、编程练习",
    },
    "expected_output": {"English": "Expected output or workload", "中文": "预期产出或工作量"},
    "expected_output_placeholder": {
        "English": "Example: 2,500 words, 6 chapters, 20 practice questions, 10-minute presentation",
        "中文": "例如：2500 字、6 个章节、20 道练习题、10 分钟展示",
    },
    "course_or_discipline": {"English": "Course or discipline", "中文": "课程或学科"},
    "course_placeholder": {
        "English": "Example: education, statistics, psychology, programming, economics",
        "中文": "例如：教育学、统计学、心理学、编程、经济学",
    },
    "required_materials": {"English": "Required materials or resources", "中文": "所需材料或资源"},
    "required_materials_placeholder": {
        "English": "Example: lecture slides, textbook chapters, dataset, readings, tutorial exercises",
        "中文": "例如：课件、教材章节、数据集、阅读材料、教程练习",
    },
    "success_criteria": {"English": "Success criteria or instructor expectations", "中文": "成功标准或教师要求"},
    "success_criteria_placeholder": {
        "English": "Example: accurate explanation, critical analysis, solved problems, clear structure, correct method",
        "中文": "例如：解释准确、批判性分析、完成题目、结构清晰、方法正确",
    },
    "current_progress": {"English": "Current progress", "中文": "当前进度"},
    "current_progress_placeholder": {
        "English": "Example: I have reviewed the first topic but cannot apply it to practice questions yet.",
        "中文": "例如：我已经复习了第一个主题，但还不能把它用到练习题中。",
    },
    "existing_materials": {"English": "Existing notes, materials, or ideas", "中文": "已有笔记、材料或想法"},
    "existing_materials_placeholder": {
        "English": "Example: I already have lecture notes, rough ideas, flashcards, or unfinished practice work.",
        "中文": "例如：我已有课堂笔记、初步想法、抽认卡或未完成的练习。",
    },
    "learning_materials": {"English": "Learning Materials", "中文": "学习材料"},
    "material_upload_intro": {
        "English": "Upload text, PDF, Word, PowerPoint, spreadsheet, or image files. PlanWise extracts available text locally and sends a short material brief to GenAI.",
        "中文": "上传文本、PDF、Word、PowerPoint、表格或图片文件。PlanWise 会在本地提取可用文本，并把简短材料摘要传给 GenAI。",
    },
    "upload_materials": {"English": "Upload materials", "中文": "上传材料"},
    "process_materials": {"English": "Process uploaded materials", "中文": "处理上传材料"},
    "materials_processed": {"English": "Learning materials processed.", "中文": "学习材料已处理。"},
    "material_processing_warning": {
        "English": "Some files could not be fully parsed. PlanWise saved what it could.",
        "中文": "部分文件无法完整解析。PlanWise 已保存可读取的内容。",
    },
    "material_library": {"English": "Material Library", "中文": "材料库"},
    "clear_materials": {"English": "Clear materials", "中文": "清空材料"},
    "materials_cleared": {"English": "Materials cleared.", "中文": "材料已清空。"},
    "image_material_note": {
        "English": "Image saved as a reference. Text interpretation requires OCR or a vision-capable model in a later version.",
        "中文": "图片已作为参考保存。图片文字识别需要后续接入 OCR 或支持视觉的模型。",
    },
    "no_material_text": {
        "English": "No extractable text was found.",
        "中文": "未提取到可用文本。",
    },
    "materials_used_in_planning": {
        "English": "Uploaded materials will be considered in Planning and AI Coach responses.",
        "中文": "已上传材料会用于 Planning 和 AI Coach 的回答。",
    },
    "generate_study_plan": {"English": "Generate Study Plan", "中文": "生成学习计划"},
    "task_required_warning": {
        "English": "Please describe your learning goal or task first.",
        "中文": "请先描述你的学习目标或任务。",
    },
    "planning_spinner": {
        "English": "Generating an SRL-informed study plan with GenAI...",
        "中文": "正在用 GenAI 生成基于自我调节学习的学习计划...",
    },
    "plan_success": {
        "English": "Study plan generated. Move to Monitoring after each study session.",
        "中文": "学习计划已生成。每次学习后，请进入监控阶段记录情况。",
    },
    "structured_plan": {"English": "Structured Study Plan", "中文": "结构化学习计划"},
    "recommended_strategies": {"English": "Recommended Learning Strategies", "中文": "推荐学习策略"},
    "generate_plan_info": {
        "English": "Generate a plan to activate the monitoring and control phases.",
        "中文": "生成计划后即可启用监控和调节阶段。",
    },
    "monitoring_dashboard": {"English": "Monitoring Dashboard", "中文": "学习监控面板"},
    "monitoring_need_plan": {
        "English": "Create a study plan in the Planning phase before monitoring progress.",
        "中文": "请先在计划阶段创建学习计划，再监控学习进度。",
    },
    "monitoring_intro": {
        "English": "After each study session, record progress signals and short evidence. These monitoring notes guide the Control phase.",
        "中文": "每次学习后，记录进度信号和简短证据。这些监控记录会帮助后续调节学习策略。",
    },
    "completed": {"English": "Completed", "中文": "已完成"},
    "delayed": {"English": "Delayed", "中文": "延迟"},
    "confused": {"English": "Confused", "中文": "困惑"},
    "overwhelmed": {"English": "Overwhelmed", "中文": "负荷过高"},
    "pending": {"English": "Pending", "中文": "未记录"},
    "multi_signal_caption": {
        "English": "Multiple signals selected. PlanWise will prioritise '{status}' for adaptive feedback.",
        "中文": "你选择了多个信号。PlanWise 会优先使用“{status}”来生成适应性反馈。",
    },
    "evidence_progress": {"English": "Evidence of progress", "中文": "进展证据"},
    "evidence_placeholder": {
        "English": "Example: solved 5 practice questions, made a concept map, drafted 2 slides, summarised one reading.",
        "中文": "例如：完成 5 道练习题、制作概念图、草拟 2 页幻灯片、总结一篇材料。",
    },
    "confidence_after_step": {"English": "Confidence after this step", "中文": "完成此步骤后的信心"},
    "barrier_point": {"English": "Barrier or confusion point", "中文": "障碍或困惑点"},
    "barrier_placeholder": {
        "English": "Example: I understand the definition but cannot apply it to a new example.",
        "中文": "例如：我理解定义，但不能把它应用到新例子中。",
    },
    "actual_time": {"English": "Actual time spent", "中文": "实际用时"},
    "tracked_completion": {"English": "Tracked completion", "中文": "已记录完成率"},
    "adaptive_feedback": {"English": "Adaptive Feedback", "中文": "适应性反馈"},
    "generate_monitoring_feedback": {"English": "Generate GenAI Monitoring Feedback", "中文": "生成 GenAI 监控反馈"},
    "monitoring_spinner": {"English": "Generating monitoring feedback with GenAI...", "中文": "正在用 GenAI 生成监控反馈..."},
    "monitoring_next_step_title": {"English": "Continue the SRL Cycle", "中文": "继续自我调节学习循环"},
    "monitoring_next_step_hint": {
        "English": "You have recorded this step. Continue to Control if you need strategy adjustment, then Reflection to complete the step and advance.",
        "中文": "你已经记录了当前步骤。若需要调整策略，请进入调节；之后进入反思，完成本步骤并推进到下一步。",
    },
    "go_to_control": {"English": "Go to Control", "中文": "进入调节"},
    "go_to_reflection": {"English": "Go to Reflection", "中文": "进入反思"},
    "back_to_monitoring": {"English": "Back to Monitoring", "中文": "返回监控"},
    "control_next_step_title": {"English": "Ready for Reflection", "中文": "准备进入反思"},
    "control_next_step_hint": {
        "English": "Use the adjustment above, then continue to Reflection to close this step and unlock the next one.",
        "中文": "参考上面的调节建议后，进入反思来完成当前步骤并解锁下一步。",
    },
    "record_current_step_first": {
        "English": "Select at least one monitoring signal before continuing.",
        "中文": "请先为当前步骤选择至少一个监控信号，再继续。",
    },
    "output_source": {"English": "Output source", "中文": "输出来源"},
    "workflow_progress": {"English": "Workflow Progress", "中文": "流程进度"},
    "current_step": {"English": "Current step", "中文": "当前步骤"},
    "step_of_total": {"English": "Step {current} of {total}", "中文": "第 {current} / {total} 步"},
    "completed_steps": {"English": "Completed steps", "中文": "已完成步骤"},
    "previous_steps": {"English": "Previous steps", "中文": "过去步骤"},
    "upcoming_steps": {"English": "Upcoming steps", "中文": "后续步骤"},
    "edit_previous_steps": {"English": "Edit previous steps", "中文": "编辑过去步骤"},
    "read_only_note": {
        "English": "Previous steps are read-only by default. Enable editing only when you need to revise them.",
        "中文": "过去步骤默认只读。只有需要修改时再开启编辑。",
    },
    "advance_next_step": {"English": "Advance to next step", "中文": "推进到下一步"},
    "workflow_complete": {"English": "Workflow complete", "中文": "流程已完成"},
    "finish_current_step_first": {
        "English": "Complete the current step reflection before advancing.",
        "中文": "请先完成当前步骤反思，再推进到下一步。",
    },
    "current_step_completed": {"English": "Current step completed.", "中文": "当前步骤已完成。"},
    "all_steps_completed": {"English": "All plan steps are completed.", "中文": "所有计划步骤已完成。"},
    "step_already_completed": {
        "English": "This step has already been completed. Return to Monitoring to continue.",
        "中文": "当前步骤已经完成。请返回监控继续下一步。",
    },
    "step_reflection_archive": {"English": "Step Reflection Archive", "中文": "步骤反思档案"},
    "completed_at_label": {"English": "Completed at", "中文": "完成时间"},
    "saved_reflection_note": {"English": "Saved reflection note", "中文": "已保存反思"},
    "no_saved_reflection_note": {"English": "No written reflection note saved.", "中文": "没有保存文字反思。"},
    "control_title": {"English": "Adaptive Strategy Adjustment", "中文": "适应性策略调节"},
    "control_need_plan": {
        "English": "Create and monitor a plan before adjusting strategies.",
        "中文": "请先创建并监控学习计划，再进行策略调节。",
    },
    "control_intro": {
        "English": "The Control phase turns monitoring signals into strategy changes, time trade-offs, and revised priorities.",
        "中文": "调节阶段会把监控信号转化为策略变化、时间取舍和优先级调整。",
    },
    "generate_control_adjustment": {"English": "Generate GenAI Plan Adjustment", "中文": "生成 GenAI 计划调节"},
    "control_spinner": {"English": "Generating control-phase adjustments with GenAI...", "中文": "正在用 GenAI 生成调节阶段建议..."},
    "updated_recommendations": {"English": "Updated Recommendations", "中文": "更新后的建议"},
    "adjusted_plan": {"English": "Adjusted Plan", "中文": "调整后的计划"},
    "monitoring_status": {"English": "Monitoring status", "中文": "监控状态"},
    "control_preferences": {"English": "Adjustment Preferences", "中文": "调节偏好"},
    "control_preferences_caption": {
        "English": "Use these optional controls to tell PlanWise how you want the plan adjusted.",
        "中文": "这些选项都是可选的，用来告诉 PlanWise 你希望如何调整计划。",
    },
    "control_goal": {"English": "Main adjustment goal", "中文": "主要调节目标"},
    "control_goal_options": {
        "English": ["Recover lost time", "Clarify understanding", "Reduce workload", "Improve quality", "Build confidence"],
        "中文": ["追回进度", "澄清理解", "降低负荷", "提升质量", "建立信心"],
    },
    "control_style": {"English": "Preferred adjustment style", "中文": "偏好的调节方式"},
    "control_style_options": {
        "English": ["Compress tasks", "Break tasks into smaller steps", "Add support activities", "Reorder priorities", "Protect revision time"],
        "中文": ["压缩任务", "拆成更小步骤", "增加支持活动", "重排优先级", "保留复习/修订时间"],
    },
    "available_time_next": {"English": "Available time for next session", "中文": "下次学习可用时间"},
    "non_negotiables": {"English": "Non-negotiables or constraints", "中文": "不能改变的要求或限制"},
    "non_negotiables_placeholder": {
        "English": "Example: I must finish the presentation slides before Friday; I cannot extend the deadline.",
        "中文": "例如：我必须在周五前完成展示幻灯片；截止日期不能延后。",
    },
    "control_risk_snapshot": {"English": "Control Risk Snapshot", "中文": "调节风险概览"},
    "risk_level": {"English": "Risk level", "中文": "风险等级"},
    "primary_signal": {"English": "Primary signal", "中文": "主要信号"},
    "low_risk": {"English": "Low", "中文": "低"},
    "medium_risk": {"English": "Medium", "中文": "中"},
    "high_risk": {"English": "High", "中文": "高"},
    "next_actions": {"English": "Next Action Protocol", "中文": "下一步行动方案"},
    "reflection_title": {"English": "Reflection Report", "中文": "反思报告"},
    "reflection_need_plan": {
        "English": "Create a plan first to generate an SRL learner profile.",
        "中文": "请先创建学习计划，再生成自我调节学习画像。",
    },
    "generate_reflection_report": {"English": "Generate GenAI Reflection Report", "中文": "生成 GenAI 反思报告"},
    "reflection_spinner": {"English": "Generating reflection report with GenAI...", "中文": "正在用 GenAI 生成反思报告..."},
    "planning_strength": {"English": "Planning strength", "中文": "计划能力"},
    "monitoring_awareness": {"English": "Monitoring awareness", "中文": "监控意识"},
    "strategy_flexibility": {"English": "Strategy flexibility", "中文": "策略灵活性"},
    "learner_profile_summary": {"English": "SRL Learner Profile Summary", "中文": "自我调节学习画像总结"},
    "reflection_prompts": {"English": "Reflection Prompts", "中文": "反思提示"},
    "reflection_notes": {"English": "Reflection notes", "中文": "反思记录"},
    "reflection_notes_placeholder": {
        "English": "Write a short reflection after today's study session.",
        "中文": "在今天的学习后写一段简短反思。",
    },
    "reflection_inputs": {"English": "Reflection Inputs", "中文": "反思输入"},
    "reflection_inputs_caption": {
        "English": "Use these optional prompts to help PlanWise distinguish effort, strategy, emotion, and next-cycle planning.",
        "中文": "这些提示都是可选的，可帮助 PlanWise 区分努力、策略、情绪和下一轮计划。",
    },
    "what_worked": {"English": "What worked well?", "中文": "哪些做法有效？"},
    "what_worked_placeholder": {
        "English": "Example: Concept mapping helped me connect ideas before practice.",
        "中文": "例如：概念图帮助我在练习前把想法联系起来。",
    },
    "what_did_not_work": {"English": "What did not work?", "中文": "哪些做法效果不好？"},
    "what_did_not_work_placeholder": {
        "English": "Example: I spent too long reviewing passively and did not test myself.",
        "中文": "例如：我花太久被动复习，没有及时自测。",
    },
    "progress_reason": {"English": "Main reason for progress or delay", "中文": "进展或延迟的主要原因"},
    "progress_reason_placeholder": {
        "English": "Example: I underestimated how hard the second topic was.",
        "中文": "例如：我低估了第二个主题的难度。",
    },
    "emotion_energy": {"English": "Emotion or energy level", "中文": "情绪或精力状态"},
    "emotion_energy_options": {
        "English": ["Focused", "Neutral", "Tired", "Anxious", "Frustrated", "Confident"],
        "中文": ["专注", "一般", "疲惫", "焦虑", "受挫", "有信心"],
    },
    "next_cycle_goal": {"English": "Next-cycle learning goal", "中文": "下一轮学习目标"},
    "next_cycle_goal_placeholder": {
        "English": "Example: Use active recall before opening my notes tomorrow.",
        "中文": "例如：明天打开笔记前先做主动回忆。",
    },
    "strategy_to_keep": {"English": "Strategy to keep", "中文": "想保留的策略"},
    "strategy_to_change": {"English": "Strategy to change", "中文": "想改变的策略"},
    "reflection_insights": {"English": "Reflection Insights", "中文": "反思洞察"},
    "next_cycle_plan": {"English": "Next-Cycle Plan", "中文": "下一轮计划"},
    "coach_panel": {"English": "AI Learning Coach Panel", "中文": "AI 学习教练面板"},
    "coach_intro": {
        "English": "Ask for metacognitive support. The coach responds with structured SRL guidance rather than open-ended chat.",
        "中文": "你可以寻求元认知支持。教练会提供结构化的自我调节学习指导，而不是开放式闲聊。",
    },
    "api_detected": {"English": "OpenAI-compatible API detected", "中文": "已检测到 OpenAI-compatible API"},
    "fallback_active": {"English": "Rule-based fallback active", "中文": "当前使用规则备用模式"},
    "example_prompts": {"English": "Example prompts", "中文": "示例问题"},
    "example_prompt_1": {"English": "I am stuck starting my study plan", "中文": "我不知道该如何开始学习计划"},
    "example_prompt_2": {
        "English": "I don't understand how to organise what I need to learn",
        "中文": "我不知道如何整理需要学习的内容",
    },
    "example_prompt_3": {"English": "I am running out of time", "中文": "我的时间快不够了"},
    "student_concern": {"English": "Student concern", "中文": "学习者困惑"},
    "student_concern_placeholder": {
        "English": "Describe where you are stuck in your learning process.",
        "中文": "描述你在学习过程中卡住的地方。",
    },
    "ask_coach": {"English": "Ask PlanWise Coach", "中文": "询问 PlanWise 教练"},
    "coach_spinner": {
        "English": "PlanWise Coach is generating a GenAI response...",
        "中文": "PlanWise 教练正在生成 GenAI 回应...",
    },
    "coach_warning": {"English": "Please enter a concern or select an example prompt.", "中文": "请输入你的困惑，或选择一个示例问题。"},
    "coaching_history": {"English": "Coaching History", "中文": "教练对话记录"},
    "student": {"English": "Student", "中文": "学习者"},
    "planwise_coach": {"English": "PlanWise Coach", "中文": "PlanWise 教练"},
    "coach_settings": {"English": "Coach Settings", "中文": "教练设置"},
    "coach_settings_caption": {
        "English": "Use these options to shape how PlanWise scaffolds your thinking.",
        "中文": "使用这些选项来决定 PlanWise 如何支架你的思考过程。",
    },
    "learning_situation": {"English": "Learning situation", "中文": "学习情境"},
    "learning_situation_options": {
        "English": ["Getting started", "Understanding concepts", "Organising materials", "Time pressure", "Motivation or anxiety", "Preparing output"],
        "中文": ["开始任务", "理解概念", "整理材料", "时间压力", "动力或焦虑", "准备产出"],
    },
    "srl_focus": {"English": "SRL focus", "中文": "自我调节学习重点"},
    "srl_focus_options": {
        "English": ["Planning", "Monitoring", "Control", "Reflection"],
        "中文": ["计划", "监控", "调节", "反思"],
    },
    "coach_response_style": {"English": "Response style", "中文": "回复风格"},
    "coach_response_style_options": {
        "English": ["Step-by-step", "Brief and direct", "Socratic questions", "Encouraging coach", "Action checklist"],
        "中文": ["分步骤指导", "简短直接", "苏格拉底式提问", "鼓励型教练", "行动清单"],
    },
    "urgency": {"English": "Urgency", "中文": "紧急程度"},
    "urgency_options": {
        "English": ["Normal", "Need help today", "Deadline is very close"],
        "中文": ["正常", "今天需要帮助", "截止日期非常近"],
    },
    "use_planner_context": {"English": "Use current planner context", "中文": "使用当前计划上下文"},
    "coach_context_note": {
        "English": "When enabled, the coach will use your plan, monitoring notes, control adjustments, and reflection inputs.",
        "中文": "开启后，教练会参考你的计划、监控记录、调节建议和反思输入。",
    },
    "clear_coach_history": {"English": "Clear coaching history", "中文": "清空教练记录"},
    "coach_history_cleared": {"English": "Coaching history cleared.", "中文": "教练记录已清空。"},
}


def initialize_state() -> None:
    defaults = {
        "language": "English",
        "tasks": load_tasks(),
        "current_task_id": "",
        "task_selectbox_id": "",
        "task_memory_selector": "",
        "pending_task_selection": None,
        "page_selector": "Dashboard",
        "pending_page_selection": None,
        "sidebar_flash": {},
        "study_plan": [],
        "plan_context": {},
        "learning_materials": [],
        "monitoring_status": {},
        "monitoring_details": {},
        "monitoring_feedback": {},
        "adjusted_plan": [],
        "control_recommendations": [],
        "control_next_actions": [],
        "control_context": {},
        "control_source": "",
        "reflection_report": {},
        "reflection_inputs": {},
        "coach_history": [],
        "workflow_current_index": 0,
        "workflow_completed_indices": [],
        "workflow_step_reflections": {},
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = copy.deepcopy(value)


def blank_learning_state() -> Dict[str, object]:
    return {
        "study_plan": [],
        "plan_context": {},
        "learning_materials": [],
        "monitoring_status": {},
        "monitoring_details": {},
        "monitoring_feedback": {},
        "adjusted_plan": [],
        "control_recommendations": [],
        "control_next_actions": [],
        "control_context": {},
        "control_source": "",
        "reflection_report": {},
        "reflection_inputs": {},
        "coach_history": [],
        "workflow_current_index": 0,
        "workflow_completed_indices": [],
        "workflow_step_reflections": {},
    }


def load_tasks() -> Dict[str, Dict[str, object]]:
    if not TASKS_FILE.exists():
        return {}
    try:
        with TASKS_FILE.open("r", encoding="utf-8") as file:
            data = json.load(file)
        return data if isinstance(data, dict) else {}
    except (OSError, json.JSONDecodeError):
        return {}


def write_tasks(tasks: Dict[str, Dict[str, object]]) -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    with TASKS_FILE.open("w", encoding="utf-8") as file:
        json.dump(tasks, file, ensure_ascii=False, indent=2)


def reset_learning_state() -> None:
    for key, value in blank_learning_state().items():
        st.session_state[key] = copy.deepcopy(value)
    st.session_state.current_task_id = ""
    st.session_state.task_selectbox_id = ""
    st.session_state.pending_task_selection = "__clear__"


def sync_task_selection(task_id: str) -> None:
    st.session_state.current_task_id = task_id
    st.session_state.task_selectbox_id = task_id
    st.session_state.pending_task_selection = task_id


def navigate_to_page(page: str) -> None:
    st.session_state.pending_page_selection = page
    st.rerun()


def task_title_from_state() -> str:
    task_description = st.session_state.plan_context.get("task_description", "")
    if task_description:
        return str(task_description)[:80]
    if st.session_state.study_plan:
        return str(st.session_state.study_plan[0].get("title", t("unsaved_task")))[:80]
    if st.session_state.coach_history:
        return str(st.session_state.coach_history[0].get("user", t("unsaved_task")))[:80]
    if st.session_state.learning_materials:
        return str(st.session_state.learning_materials[0].get("name", t("unsaved_task")))[:80]
    return t("unsaved_task")


def task_snapshot() -> Dict[str, object]:
    return {
        key: copy.deepcopy(st.session_state.get(key, blank_learning_state().get(key)))
        for key in PERSISTED_TASK_KEYS
    }


def save_current_task(create_if_missing: bool = False, custom_title: str | None = None) -> str:
    has_content = bool(
        st.session_state.plan_context
        or st.session_state.study_plan
        or st.session_state.coach_history
        or st.session_state.learning_materials
    )
    if not has_content:
        return ""

    task_id = st.session_state.current_task_id
    if not task_id and not create_if_missing:
        return ""
    created_new_task = False
    if not task_id:
        task_id = uuid.uuid4().hex
        created_new_task = True
        sync_task_selection(task_id)

    now = datetime.now().isoformat(timespec="seconds")
    existing = st.session_state.tasks.get(task_id, {})
    clean_custom_title = str(custom_title or "").strip()
    if clean_custom_title:
        title = clean_custom_title[:120]
        custom_title_enabled = True
    elif existing.get("custom_title"):
        title = str(existing.get("title") or task_title_from_state())
        custom_title_enabled = True
    else:
        title = task_title_from_state()
        custom_title_enabled = False
    st.session_state.tasks[task_id] = {
        "id": task_id,
        "title": title,
        "custom_title": custom_title_enabled,
        "created_at": existing.get("created_at", now),
        "updated_at": now,
        "data": task_snapshot(),
    }
    write_tasks(st.session_state.tasks)
    selected_task_id = st.session_state.get("task_selectbox_id", "")
    if created_new_task or not selected_task_id or selected_task_id == task_id:
        sync_task_selection(task_id)
    return task_id


def load_task_into_state(task_id: str) -> None:
    task = st.session_state.tasks.get(task_id)
    if not task:
        return
    reset_learning_state()
    for key, value in task.get("data", {}).items():
        if key in PERSISTED_TASK_KEYS:
            st.session_state[key] = copy.deepcopy(value)
    sync_task_selection(task_id)
    clamp_workflow_state()


def clamp_workflow_state() -> None:
    total = len(st.session_state.get("study_plan", []))
    if total <= 0:
        st.session_state.workflow_current_index = 0
        st.session_state.workflow_completed_indices = []
        return
    current = int(st.session_state.get("workflow_current_index", 0) or 0)
    st.session_state.workflow_current_index = min(max(current, 0), total - 1)
    completed = st.session_state.get("workflow_completed_indices", []) or []
    st.session_state.workflow_completed_indices = sorted(
        {
            int(index)
            for index in completed
            if isinstance(index, int) or str(index).isdigit()
        }
        & set(range(total))
    )


def current_step_index() -> int:
    clamp_workflow_state()
    return int(st.session_state.workflow_current_index)


def current_step() -> Dict[str, str] | None:
    plan = st.session_state.get("study_plan", [])
    if not plan:
        return None
    return plan[current_step_index()]


def is_current_step_recorded() -> bool:
    step = current_step()
    if not step:
        return False
    day = step.get("day", "")
    status = st.session_state.monitoring_status.get(day, "Pending")
    details = st.session_state.monitoring_details.get(day, {})
    has_note = bool(str(details.get("evidence", "")).strip() or str(details.get("barrier", "")).strip())
    return status != "Pending" or has_note


def mark_current_step_complete() -> None:
    clamp_workflow_state()
    total = len(st.session_state.study_plan)
    if total <= 0:
        return
    current = current_step_index()
    if current in set(st.session_state.workflow_completed_indices):
        return
    reflections = st.session_state.get("workflow_step_reflections", {}) or {}
    reflections[str(current)] = {
        "step": copy.deepcopy(current_step()),
        "reflection_inputs": copy.deepcopy(st.session_state.get("reflection_inputs", {})),
        "reflection_report": copy.deepcopy(st.session_state.get("reflection_report", {})),
        "completed_at": datetime.now().isoformat(timespec="seconds"),
    }
    st.session_state.workflow_step_reflections = reflections

    completed = set(st.session_state.workflow_completed_indices)
    completed.add(current)
    st.session_state.workflow_completed_indices = sorted(completed)
    if current < total - 1:
        st.session_state.workflow_current_index = current + 1
        st.session_state.monitoring_feedback = {}
        st.session_state.control_recommendations = []
        st.session_state.control_next_actions = []
        st.session_state.control_context = {}
        st.session_state.control_source = ""
        st.session_state.reflection_report = {}
        st.session_state.reflection_inputs = {}
    save_current_task(create_if_missing=True)


def render_workflow_progress() -> None:
    plan = st.session_state.get("study_plan", [])
    if not plan:
        return
    clamp_workflow_state()
    total = len(plan)
    current = current_step_index()
    completed_count = len(st.session_state.workflow_completed_indices)
    st.subheader(t("workflow_progress"))
    st.progress(completed_count / total if total else 0)
    col1, col2, col3 = st.columns(3)
    col1.metric(t("current_step"), t("step_of_total").format(current=current + 1, total=total))
    col2.metric(t("completed_steps"), f"{completed_count}/{total}")
    col3.metric(t("tracked_completion"), f"{round((completed_count / total) * 100, 1)}%")


def render_step_summary(step: Dict[str, str], index: int) -> None:
    status = st.session_state.monitoring_status.get(step.get("day", ""), "Pending")
    details = st.session_state.monitoring_details.get(step.get("day", ""), {})
    evidence = details.get("evidence", "")
    barrier = details.get("barrier", "")
    marker = "✓ " if index in st.session_state.workflow_completed_indices else ""
    st.markdown(
        f"""
        <div class="plan-day">
            <strong>{marker}{escape(str(step.get("day", "")))}: {escape(str(step.get("title", "")))}</strong><br>
            <span>{escape(str(step.get("detail", "")))}</span><br>
            <small>{t("monitoring_status")}: {status_label(status)}</small>
            {f'<br><small>{escape(str(evidence))}</small>' if evidence else ''}
            {f'<br><small>{escape(str(barrier))}</small>' if barrier else ''}
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_step_reflection_archive() -> None:
    reflections = st.session_state.get("workflow_step_reflections", {}) or {}
    if not reflections:
        return

    def archive_sort_key(item: Tuple[str, object]) -> int:
        key, _ = item
        return int(key) if str(key).isdigit() else 0

    with st.expander(t("step_reflection_archive"), expanded=False):
        for index_text, archive in sorted(reflections.items(), key=archive_sort_key):
            if not isinstance(archive, dict):
                continue
            step = archive.get("step") or {}
            inputs = archive.get("reflection_inputs") or {}
            completed_at = archive.get("completed_at", "")
            note = ""
            if isinstance(inputs, dict):
                note = str(inputs.get("reflection_notes", "")).strip()
            step_number = int(index_text) + 1 if str(index_text).isdigit() else index_text
            title = step.get("title", f"{t('current_step')} {step_number}") if isinstance(step, dict) else ""
            day = step.get("day", "") if isinstance(step, dict) else ""
            st.markdown(
                f"""
                <div class="plan-day">
                    <strong>{t("step_of_total").format(current=step_number, total=len(st.session_state.study_plan))}: {escape(str(day))} {escape(str(title))}</strong><br>
                    <small>{t("completed_at_label")}: {escape(str(completed_at))}</small><br>
                    <span><strong>{t("saved_reflection_note")}:</strong> {escape(note or t("no_saved_reflection_note"))}</span>
                </div>
                """,
                unsafe_allow_html=True,
            )


def safe_material_filename(filename: str) -> str:
    name = Path(filename).name
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("._")
    return cleaned or f"material_{uuid.uuid4().hex}"


def clip_text(text: str, limit: int = 6000) -> str:
    cleaned = re.sub(r"\s+", " ", text or "").strip()
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[:limit].rstrip() + "..."


def extract_text_from_text_file(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise RuntimeError("Install pypdf to parse PDF files.") from exc

    reader = PdfReader(io.BytesIO(file_bytes))
    page_text = []
    for page_index, page in enumerate(reader.pages[:40], start=1):
        text = page.extract_text() or ""
        if text.strip():
            page_text.append(f"[Page {page_index}] {text}")
    return "\n\n".join(page_text)


def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
    except Exception as exc:
        raise RuntimeError("Install python-docx to parse DOCX files.") from exc

    document = Document(io.BytesIO(file_bytes))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("Install python-pptx to parse PPTX files.") from exc

    presentation = Presentation(io.BytesIO(file_bytes))
    slides = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
        if parts:
            slides.append(f"[Slide {slide_index}] " + "\n".join(parts))
    return "\n\n".join(slides)


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise RuntimeError("Install openpyxl to parse XLSX files.") from exc

    workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets[:5]:
        parts.append(f"[Sheet: {sheet.title}]")
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > 80:
                parts.append("[Rows truncated]")
                break
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_from_csv(file_bytes: bytes) -> str:
    text = extract_text_from_text_file(file_bytes)
    rows = []
    reader = csv.reader(io.StringIO(text))
    for row_index, row in enumerate(reader, start=1):
        if row_index > 120:
            rows.append("[Rows truncated]")
            break
        cleaned = [cell.strip() for cell in row if cell.strip()]
        if cleaned:
            rows.append(" | ".join(cleaned))
    return "\n".join(rows)


def image_metadata(file_bytes: bytes) -> Dict[str, object]:
    try:
        from PIL import Image
    except Exception as exc:
        raise RuntimeError("Install Pillow to inspect image files.") from exc

    with Image.open(io.BytesIO(file_bytes)) as image:
        return {
            "width": image.width,
            "height": image.height,
            "mode": image.mode,
            "format": image.format,
        }


def parse_uploaded_material(uploaded_file) -> Dict[str, object]:
    file_bytes = uploaded_file.getvalue()
    filename = uploaded_file.name
    extension = Path(filename).suffix.lower().lstrip(".")
    material_id = uuid.uuid4().hex
    safe_name = f"{material_id}_{safe_material_filename(filename)}"
    saved_path = UPLOAD_DIR / safe_name
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    saved_path.write_bytes(file_bytes)

    material = {
        "id": material_id,
        "name": filename,
        "extension": extension,
        "size_bytes": len(file_bytes),
        "saved_path": str(saved_path),
        "uploaded_at": datetime.now().isoformat(timespec="seconds"),
        "text": "",
        "excerpt": "",
        "status": "processed",
        "note": "",
        "metadata": {},
    }

    try:
        if extension in {"txt", "md"}:
            text = extract_text_from_text_file(file_bytes)
        elif extension == "pdf":
            text = extract_text_from_pdf(file_bytes)
        elif extension == "docx":
            text = extract_text_from_docx(file_bytes)
        elif extension == "pptx":
            text = extract_text_from_pptx(file_bytes)
        elif extension == "xlsx":
            text = extract_text_from_xlsx(file_bytes)
        elif extension == "csv":
            text = extract_text_from_csv(file_bytes)
        elif extension in {"png", "jpg", "jpeg", "webp"}:
            material["metadata"] = image_metadata(file_bytes)
            text = ""
            material["note"] = t("image_material_note")
        else:
            text = ""
            material["status"] = "unsupported"
            material["note"] = f"Unsupported file type: .{extension}"

        material["text"] = clip_text(text, 16000)
        material["excerpt"] = clip_text(text, 1200)
        if not material["text"] and not material["note"]:
            material["note"] = t("no_material_text")
    except Exception as exc:
        material["status"] = "partial"
        material["note"] = str(exc)

    return material


def materials_brief(limit: int = 6000) -> str:
    materials = st.session_state.get("learning_materials", []) or []
    if not materials:
        return ""
    parts = []
    for index, material in enumerate(materials, start=1):
        name = material.get("name", f"Material {index}")
        extension = material.get("extension", "")
        note = material.get("note", "")
        excerpt = material.get("excerpt", "") or material.get("text", "")
        metadata = material.get("metadata", {})
        descriptor = f"{index}. {name} ({extension})"
        if metadata:
            descriptor += f" metadata={metadata}"
        if excerpt:
            descriptor += f"\nExcerpt: {excerpt}"
        elif note:
            descriptor += f"\nNote: {note}"
        parts.append(descriptor)
    return clip_text("\n\n".join(parts), limit)


def render_material_library() -> None:
    materials = st.session_state.get("learning_materials", []) or []
    if not materials:
        return
    st.subheader(t("material_library"))
    st.caption(t("materials_used_in_planning"))
    for material in materials:
        status = escape(str(material.get("status", "")))
        name = escape(str(material.get("name", "")))
        note = escape(str(material.get("note", "")))
        excerpt = escape(str(material.get("excerpt", "")) or note or t("no_material_text"))
        st.markdown(
            f"""
            <div class="section-band">
                <strong>{name}</strong><br>
                <small>{status} · {escape(str(material.get("size_bytes", 0)))} bytes</small><br>
                <span>{excerpt}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
    if st.button(t("clear_materials")):
        st.session_state.learning_materials = []
        save_current_task()
        st.success(t("materials_cleared"))
        st.rerun()


def render_material_upload_panel(key_prefix: str) -> None:
    with st.expander(t("learning_materials"), expanded=False):
        st.caption(t("material_upload_intro"))
        uploaded_files = st.file_uploader(
            t("upload_materials"),
            type=SUPPORTED_MATERIAL_TYPES,
            accept_multiple_files=True,
            key=f"{key_prefix}_material_uploader",
        )
        if st.button(t("process_materials"), key=f"{key_prefix}_process_materials"):
            if uploaded_files:
                parsed_materials = [parse_uploaded_material(file) for file in uploaded_files]
                existing_ids = {material.get("id") for material in st.session_state.learning_materials}
                for material in parsed_materials:
                    if material.get("id") not in existing_ids:
                        st.session_state.learning_materials.append(material)
                save_current_task(create_if_missing=True)
                if any(material.get("status") != "processed" for material in parsed_materials):
                    st.warning(t("material_processing_warning"))
                else:
                    st.success(t("materials_processed"))
                st.rerun()
        render_material_library()


def rename_task(task_id: str, new_title: str) -> bool:
    new_title = new_title.strip()
    if not task_id or not new_title or task_id not in st.session_state.tasks:
        return False
    st.session_state.tasks[task_id]["title"] = new_title[:120]
    st.session_state.tasks[task_id]["custom_title"] = True
    st.session_state.tasks[task_id]["updated_at"] = datetime.now().isoformat(timespec="seconds")
    write_tasks(st.session_state.tasks)
    st.session_state.task_selectbox_id = task_id
    st.session_state.pending_task_selection = task_id
    return True


def delete_task(task_id: str) -> bool:
    if not task_id or task_id not in st.session_state.tasks:
        return False
    del st.session_state.tasks[task_id]
    write_tasks(st.session_state.tasks)
    if st.session_state.current_task_id == task_id:
        reset_learning_state()
    elif st.session_state.task_selectbox_id == task_id:
        st.session_state.task_selectbox_id = ""
        st.session_state.pending_task_selection = "__clear__"
    return True


def task_preview_text(task_id: str) -> str:
    task = st.session_state.tasks.get(task_id, {})
    data = task.get("data", {})
    plan_context = data.get("plan_context", {})
    study_plan = data.get("study_plan", [])
    learning_materials = data.get("learning_materials", [])
    deadline = plan_context.get("deadline", "Not set")
    task_description = plan_context.get("task_description", task.get("title", ""))
    plan_count = len(study_plan) if isinstance(study_plan, list) else 0
    material_count = len(learning_materials) if isinstance(learning_materials, list) else 0
    updated_at = str(task.get("updated_at", "")).replace("T", " ")
    created_at = str(task.get("created_at", "")).replace("T", " ")
    if current_language() == "中文":
        return (
            f"任务：{task_description}\n"
            f"截止日期：{deadline}\n"
            f"计划项：{plan_count}\n"
            f"材料数：{material_count}\n"
            f"创建时间：{created_at}\n"
            f"更新时间：{updated_at}"
        )
    return (
        f"Task: {task_description}\n"
        f"Deadline: {deadline}\n"
        f"Plan items: {plan_count}\n"
        f"Materials: {material_count}\n"
        f"Created: {created_at}\n"
        f"Updated: {updated_at}"
    )


def current_language() -> str:
    return st.session_state.get("language", "English")


def t(key: str) -> str:
    value = TRANSLATIONS.get(key, {})
    if isinstance(value, dict):
        return value.get(current_language(), value.get("English", key))
    return str(value or key)


def phase_label(phase: str) -> str:
    return {
        "Dashboard": t("phase_dashboard"),
        "Task Library": t("phase_task_library"),
        "Task Details": t("phase_task_details"),
        "Planning": t("phase_planning"),
        "Today Learning": t("phase_today"),
        "Monitoring": t("phase_monitoring"),
        "Control": t("phase_control"),
        "Reflection": t("phase_reflection"),
        "AI Coach": t("phase_ai_coach"),
        "Reports": t("phase_reports"),
        "How to Use": t("phase_guide"),
    }.get(phase, phase)


def status_label(status: str) -> str:
    return {
        "Pending": t("pending"),
        "Completed": t("completed"),
        "Delayed": t("delayed"),
        "Confused": t("confused"),
        "Overwhelmed": t("overwhelmed"),
    }.get(status, status)


def difficulty_label(difficulty: str) -> str:
    if current_language() == "中文":
        return {"Low": "低", "Medium": "中", "High": "高"}.get(difficulty, difficulty)
    return difficulty


def level_label(level: object) -> str:
    text = str(level)
    if current_language() == "中文":
        return {"High": "高", "Medium": "中", "Developing": "发展中"}.get(text, text)
    return text


def output_language_name() -> str:
    return "Simplified Chinese" if current_language() == "中文" else "English"


def translated_source(source: str) -> str:
    if current_language() != "中文":
        return source
    if "GenAI-generated" in source:
        return "由 GenAI 通过 API 生成"
    if "Rule-based preview" in source:
        return "GenAI 生成前的规则预览"
    if "Rule-based fallback" in source:
        return "GenAI 不可用后的规则备用输出"
    return source


def task_option_label(task_id: str) -> str:
    if not task_id:
        return t("unsaved_task")
    task = st.session_state.tasks.get(task_id, {})
    title = task.get("title", t("unsaved_task"))
    updated_at = str(task.get("updated_at", ""))[:16].replace("T", " ")
    return f"{title} · {updated_at}" if updated_at else str(title)


def set_sidebar_flash(message_key: str, level: str = "success") -> None:
    st.session_state.sidebar_flash = {"message_key": message_key, "level": level}


def render_sidebar_flash() -> None:
    flash = st.session_state.get("sidebar_flash", {}) or {}
    if not flash:
        return
    message = t(str(flash.get("message_key", "")))
    level = str(flash.get("level", "success"))
    if level == "warning":
        st.sidebar.warning(message)
    elif level == "error":
        st.sidebar.error(message)
    else:
        st.sidebar.success(message)
    st.session_state.sidebar_flash = {}


def inject_styles() -> None:
    st.markdown(
        """
        <style>
        :root {
            --planwise-ink: #18212f;
            --planwise-muted: #667085;
            --planwise-line: #d8e0ea;
            --planwise-sky: #e7f1ff;
            --planwise-teal: #ddf7ee;
            --planwise-gold: #fff2cc;
            --planwise-rose: #ffe5e5;
            --planwise-green: #207a63;
            --planwise-blue: #315fdc;
            --planwise-coral: #f26d5b;
            --planwise-violet: #7257d6;
            --planwise-bg: #f5f7fb;
            --planwise-card: #ffffff;
            --planwise-card-soft: #fbfdff;
            --planwise-field-bg: #ffffff;
            --planwise-field-text: #18212f;
            --planwise-hero-bg: linear-gradient(135deg, #eef5ff 0%, #f6fbf2 54%, #fff2ee 100%);
            --planwise-app-bg: radial-gradient(circle at 8% 4%, rgba(49, 95, 220, 0.13), transparent 28rem),
                linear-gradient(180deg, #f8fbff 0%, var(--planwise-bg) 55%, #ffffff 100%);
            --planwise-shadow: 0 8px 22px rgba(24, 33, 47, 0.05);
        }
        [data-testid="stAppViewContainer"] {
            background: var(--planwise-app-bg);
            color: var(--planwise-ink);
        }
        [data-testid="stSidebar"] {
            background: #111827;
        }
        [data-testid="stSidebar"] * {
            color: #f9fafb;
        }
        [data-testid="stSidebar"] input,
        [data-testid="stSidebar"] textarea,
        [data-testid="stSidebar"] select {
            color: #f9fafb !important;
        }
        .main .block-container {
            padding-top: 1.45rem;
            max-width: 1220px;
        }
        h1, h2, h3 {
            letter-spacing: 0;
            color: var(--planwise-ink) !important;
        }
        [data-testid="stAppViewContainer"] p,
        [data-testid="stAppViewContainer"] li,
        [data-testid="stAppViewContainer"] label,
        [data-testid="stAppViewContainer"] [data-testid="stMarkdownContainer"],
        [data-testid="stAppViewContainer"] [data-testid="stMarkdownContainer"] * {
            color: var(--planwise-ink);
        }
        [data-testid="stAppViewContainer"] [data-testid="stCaptionContainer"],
        [data-testid="stAppViewContainer"] [data-testid="stCaptionContainer"] * {
            color: var(--planwise-muted) !important;
        }
        [data-testid="stAppViewContainer"] [data-testid="stWidgetLabel"] *,
        [data-testid="stAppViewContainer"] [data-testid="stForm"] label *,
        [data-testid="stAppViewContainer"] [data-testid="stExpander"] summary * {
            color: var(--planwise-ink) !important;
        }
        [data-testid="stAppViewContainer"] input,
        [data-testid="stAppViewContainer"] textarea,
        [data-testid="stAppViewContainer"] select,
        [data-testid="stAppViewContainer"] [role="combobox"],
        [data-testid="stAppViewContainer"] [data-baseweb="select"] > div,
        [data-testid="stAppViewContainer"] [data-baseweb="input"] > div,
        [data-testid="stAppViewContainer"] [data-baseweb="textarea"] > div {
            background-color: var(--planwise-field-bg) !important;
            color: var(--planwise-field-text) !important;
            border-color: var(--planwise-line) !important;
        }
        [data-testid="stAppViewContainer"] input::placeholder,
        [data-testid="stAppViewContainer"] textarea::placeholder {
            color: var(--planwise-muted) !important;
            opacity: 1 !important;
        }
        [data-testid="stAppViewContainer"] [data-baseweb="tag"] {
            background-color: var(--planwise-coral) !important;
            color: #ffffff !important;
            border-radius: 8px !important;
            margin: 0.25rem 0.35rem 0.25rem 0.35rem !important;
            padding: 0.2rem 0.55rem !important;
            max-width: none !important;
            width: auto !important;
            min-width: fit-content !important;
            transform: none !important;
        }
        [data-testid="stAppViewContainer"] [data-baseweb="tag"] span,
        [data-testid="stAppViewContainer"] [data-baseweb="tag"] div,
        [data-testid="stAppViewContainer"] [data-baseweb="tag"] > div,
        [data-testid="stAppViewContainer"] [data-baseweb="tag"] [title] {
            color: #ffffff !important;
            max-width: none !important;
            width: auto !important;
            min-width: fit-content !important;
            overflow: visible !important;
            text-overflow: clip !important;
            white-space: nowrap !important;
            transform: none !important;
        }
        [data-testid="stAppViewContainer"] [data-baseweb="tag"] svg {
            color: #ffffff !important;
            fill: #ffffff !important;
        }
        [data-testid="stAppViewContainer"] [data-baseweb="select"] [role="listbox"] {
            padding-left: 0.75rem !important;
            gap: 0.25rem !important;
            overflow: visible !important;
        }
        [data-testid="stAppViewContainer"] [data-baseweb="select"] [role="listbox"] > div {
            overflow: visible !important;
            max-width: none !important;
        }
        .app-hero {
            background: var(--planwise-hero-bg);
            border: 1px solid rgba(49, 95, 220, 0.16);
            border-radius: 8px;
            padding: 1.15rem 1.25rem;
            margin-bottom: 1rem;
            color: var(--planwise-ink);
        }
        .app-hero h1 {
            margin: 0;
            font-size: 2.15rem;
        }
        .app-hero p {
            margin: 0.45rem 0 0 0;
            color: var(--planwise-muted);
            max-width: 760px;
        }
        .dashboard-grid {
            display: grid;
            grid-template-columns: repeat(4, minmax(0, 1fr));
            gap: 0.75rem;
            margin: 1rem 0;
        }
        .metric-card {
            background: var(--planwise-card);
            border: 1px solid var(--planwise-line);
            border-radius: 8px;
            padding: 1rem;
            color: var(--planwise-ink);
            box-shadow: var(--planwise-shadow);
            min-height: 112px;
        }
        .metric-card small {
            color: var(--planwise-muted);
            font-weight: 700;
            text-transform: uppercase;
            font-size: 0.72rem;
        }
        .metric-card strong {
            display: block;
            font-size: 1.8rem;
            margin: 0.25rem 0;
            color: var(--planwise-ink);
        }
        .task-card {
            background: var(--planwise-card);
            border: 1px solid var(--planwise-line);
            border-radius: 8px;
            padding: 1rem;
            color: var(--planwise-ink);
            margin-bottom: 0.75rem;
            box-shadow: var(--planwise-shadow);
        }
        .task-card strong {
            color: var(--planwise-ink);
        }
        .task-card span,
        .task-card small {
            color: var(--planwise-muted);
        }
        .timeline-row {
            display: grid;
            grid-template-columns: 2rem 1fr;
            gap: 0.75rem;
            align-items: start;
            margin: 0.55rem 0;
        }
        .timeline-dot {
            width: 1.3rem;
            height: 1.3rem;
            border-radius: 999px;
            background: var(--planwise-line);
            border: 3px solid var(--planwise-card);
            box-shadow: 0 0 0 1px var(--planwise-line);
            margin-top: 0.25rem;
        }
        .timeline-dot.done {
            background: var(--planwise-green);
            box-shadow: 0 0 0 1px #82d7bb;
        }
        .timeline-dot.active {
            background: var(--planwise-coral);
            box-shadow: 0 0 0 1px #f5a196;
        }
        .cycle-wrap {
            display: grid;
            grid-template-columns: repeat(4, minmax(0, 1fr));
            gap: 0.75rem;
            margin: 1rem 0 1.4rem 0;
        }
        .cycle-step {
            border: 1px solid var(--planwise-line);
            border-radius: 8px;
            padding: 0.85rem;
            background: var(--planwise-card);
            color: var(--planwise-ink);
            min-height: 92px;
        }
        .cycle-step strong {
            color: var(--planwise-ink);
            display: block;
            margin-bottom: 0.25rem;
        }
        .cycle-step span {
            color: var(--planwise-muted);
            font-size: 0.9rem;
        }
        .active-phase {
            border: 2px solid var(--planwise-blue);
            background: var(--planwise-sky);
        }
        .phase-label {
            color: var(--planwise-blue);
            font-weight: 700;
            text-transform: uppercase;
            font-size: 0.78rem;
            letter-spacing: 0.08rem;
            margin-bottom: 0.2rem;
        }
        .section-band {
            border: 1px solid var(--planwise-line);
            border-radius: 8px;
            padding: 1rem 1.1rem;
            background: var(--planwise-card-soft);
            color: var(--planwise-ink);
            margin-bottom: 1rem;
        }
        .plan-day {
            border-left: 4px solid var(--planwise-blue);
            background: var(--planwise-card);
            color: var(--planwise-ink);
            padding: 0.8rem 1rem;
            margin-bottom: 0.65rem;
            border-radius: 0 8px 8px 0;
            border-top: 1px solid var(--planwise-line);
            border-right: 1px solid var(--planwise-line);
            border-bottom: 1px solid var(--planwise-line);
            box-shadow: var(--planwise-shadow);
        }
        .plan-day strong {
            color: var(--planwise-ink);
        }
        .plan-day span {
            color: var(--planwise-muted);
        }
        .plan-day small {
            color: var(--planwise-muted);
        }
        .feedback-box {
            border-radius: 8px;
            padding: 0.9rem 1rem;
            background: var(--planwise-teal);
            border: 1px solid #b9e7dc;
            color: var(--planwise-ink);
            margin: 0.65rem 0;
        }
        .warning-box {
            border-radius: 8px;
            padding: 0.9rem 1rem;
            background: var(--planwise-gold);
            border: 1px solid #ead18e;
            color: var(--planwise-ink);
            margin: 0.65rem 0;
        }
        .coach-box {
            border-radius: 8px;
            padding: 1rem;
            background: var(--planwise-card);
            border: 1px solid var(--planwise-line);
            color: var(--planwise-ink);
            margin-bottom: 0.75rem;
        }
        .coach-box strong {
            color: var(--planwise-ink);
        }
        [data-testid="stAppViewContainer"] .plan-day,
        [data-testid="stAppViewContainer"] .cycle-step,
        [data-testid="stAppViewContainer"] .section-band,
        [data-testid="stAppViewContainer"] .coach-box,
        [data-testid="stAppViewContainer"] .metric-card,
        [data-testid="stAppViewContainer"] .task-card,
        [data-testid="stAppViewContainer"] .app-hero {
            color: var(--planwise-ink) !important;
        }
        [data-testid="stAppViewContainer"] .plan-day *,
        [data-testid="stAppViewContainer"] .coach-box * {
            color: inherit;
        }
        [data-testid="stAppViewContainer"] .plan-day span {
            color: var(--planwise-muted) !important;
        }
        [data-testid="stAppViewContainer"] .plan-day small {
            color: var(--planwise-muted) !important;
        }
        @media (prefers-color-scheme: dark) {
            :root {
                --planwise-ink: #f4f7fb;
                --planwise-muted: #c3cad7;
                --planwise-line: #344156;
                --planwise-sky: #172b4f;
                --planwise-teal: #103c34;
                --planwise-gold: #3f3313;
                --planwise-rose: #431d25;
                --planwise-green: #58d5ac;
                --planwise-blue: #8fb2ff;
                --planwise-coral: #ff8d78;
                --planwise-violet: #b9a7ff;
                --planwise-bg: #0d1320;
                --planwise-card: #151d2b;
                --planwise-card-soft: #111827;
                --planwise-field-bg: #0f1724;
                --planwise-field-text: #f4f7fb;
                --planwise-hero-bg: linear-gradient(135deg, #152641 0%, #102d2a 52%, #331d25 100%);
                --planwise-app-bg: radial-gradient(circle at 8% 4%, rgba(143, 178, 255, 0.14), transparent 28rem),
                    linear-gradient(180deg, #0d1320 0%, #111827 58%, #0b1020 100%);
                --planwise-shadow: 0 10px 26px rgba(0, 0, 0, 0.28);
            }
            [data-testid="stAppViewContainer"] {
                background: var(--planwise-app-bg) !important;
                color: var(--planwise-ink) !important;
            }
            [data-testid="stAppViewContainer"] h1,
            [data-testid="stAppViewContainer"] h2,
            [data-testid="stAppViewContainer"] h3,
            [data-testid="stAppViewContainer"] p,
            [data-testid="stAppViewContainer"] li,
            [data-testid="stAppViewContainer"] label,
            [data-testid="stAppViewContainer"] [data-testid="stMarkdownContainer"],
            [data-testid="stAppViewContainer"] [data-testid="stMarkdownContainer"] *,
            [data-testid="stAppViewContainer"] [data-testid="stWidgetLabel"] *,
            [data-testid="stAppViewContainer"] [data-testid="stForm"] label *,
            [data-testid="stAppViewContainer"] [data-testid="stExpander"] summary * {
                color: var(--planwise-ink) !important;
            }
            [data-testid="stAppViewContainer"] [data-testid="stCaptionContainer"],
            [data-testid="stAppViewContainer"] [data-testid="stCaptionContainer"] *,
            [data-testid="stAppViewContainer"] small {
                color: var(--planwise-muted) !important;
            }
            [data-testid="stAppViewContainer"] .plan-day,
            [data-testid="stAppViewContainer"] .cycle-step,
            [data-testid="stAppViewContainer"] .section-band,
            [data-testid="stAppViewContainer"] .coach-box,
            [data-testid="stAppViewContainer"] .metric-card,
            [data-testid="stAppViewContainer"] .task-card,
            [data-testid="stAppViewContainer"] .app-hero {
                background-color: var(--planwise-card) !important;
                border-color: var(--planwise-line) !important;
                color: var(--planwise-ink) !important;
            }
            [data-testid="stAppViewContainer"] input,
            [data-testid="stAppViewContainer"] textarea,
            [data-testid="stAppViewContainer"] select,
            [data-testid="stAppViewContainer"] [role="combobox"],
            [data-testid="stAppViewContainer"] [data-baseweb="select"] > div,
            [data-testid="stAppViewContainer"] [data-baseweb="input"] > div,
            [data-testid="stAppViewContainer"] [data-baseweb="textarea"] > div {
                background-color: var(--planwise-field-bg) !important;
                color: var(--planwise-field-text) !important;
                border-color: var(--planwise-line) !important;
            }
            [data-testid="stAppViewContainer"] [data-baseweb="tag"],
            [data-testid="stAppViewContainer"] [data-baseweb="tag"] span,
            [data-testid="stAppViewContainer"] [data-baseweb="tag"] div,
            [data-testid="stAppViewContainer"] [data-baseweb="tag"] [title] {
                color: #ffffff !important;
                overflow: visible !important;
                max-width: none !important;
                width: auto !important;
                min-width: fit-content !important;
            }
        }
        @media (max-width: 760px) {
            .cycle-wrap,
            .dashboard-grid {
                grid-template-columns: 1fr;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_header(active_phase: str) -> None:
    st.title("PlanWise")
    st.caption(t("app_caption"))
    descriptions = {
        "Planning": t("phase_desc_planning"),
        "Monitoring": t("phase_desc_monitoring"),
        "Control": t("phase_desc_control"),
        "Reflection": t("phase_desc_reflection"),
    }
    html = '<div class="cycle-wrap">'
    for phase in SRL_PHASES:
        active_class = " active-phase" if phase == active_phase else ""
        html += (
            f'<div class="cycle-step{active_class}">'
            f"<strong>{phase_label(phase)}</strong><span>{descriptions[phase]}</span>"
            "</div>"
        )
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


def render_hero(title: str, subtitle: str) -> None:
    st.markdown(
        f"""
        <div class="app-hero">
            <h1>{escape(title)}</h1>
            <p>{escape(subtitle)}</p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_metric_card(label: str, value: object, caption: str = "") -> str:
    return (
        f'<div class="metric-card"><small>{escape(label)}</small>'
        f'<strong>{escape(str(value))}</strong>'
        f'<span>{escape(caption)}</span></div>'
    )


def render_dashboard_metrics() -> None:
    total_steps = len(st.session_state.study_plan)
    completed_steps = len(st.session_state.workflow_completed_indices)
    progress = f"{round((completed_steps / total_steps) * 100, 1)}%" if total_steps else "0%"
    material_count = len(st.session_state.learning_materials)
    reflection_count = len(st.session_state.workflow_step_reflections)
    current_number = current_step_index() + 1 if total_steps else 0
    cards = [
        render_metric_card(t("task_progress"), progress, f"{completed_steps}/{total_steps}"),
        render_metric_card(t("current_step"), current_number, t("step_of_total").format(current=current_number, total=max(total_steps, 1))),
        render_metric_card(t("materials_count"), material_count, t("learning_materials")),
        render_metric_card(t("reflection_count"), reflection_count, t("phase_reflection")),
    ]
    st.markdown(f'<div class="dashboard-grid">{"".join(cards)}</div>', unsafe_allow_html=True)


def render_current_task_card() -> None:
    task = st.session_state.plan_context.get("task_description", "")
    deadline = st.session_state.plan_context.get("deadline", "Not set")
    difficulty = st.session_state.plan_context.get("difficulty", "")
    challenge = st.session_state.plan_context.get("challenge", "")
    if not task and not st.session_state.learning_materials:
        st.info(t("no_active_task_dashboard"))
        return
    title = task or task_title_from_state()
    st.markdown(
        f"""
        <div class="task-card">
            <small>{t("task_details")}</small><br>
            <strong>{escape(str(title))}</strong><br>
            <span>{t("deadline")}: {escape(str(deadline))}</span><br>
            <span>{t("perceived_difficulty")}: {escape(str(difficulty_label(str(difficulty))))}</span>
            {f'<br><span>{t("current_challenge")}: {escape(str(challenge))}</span>' if challenge else ''}
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_current_step_panel() -> None:
    step = current_step()
    if not step:
        return
    st.subheader(t("current_workflow_step"))
    render_step_summary(step, current_step_index())


def render_progress_timeline(limit: int | None = None) -> None:
    plan = st.session_state.study_plan
    if not plan:
        return
    st.subheader(t("progress_timeline"))
    current = current_step_index()
    completed = set(st.session_state.workflow_completed_indices)
    visible_plan = plan[:limit] if limit else plan
    html = ""
    for index, step in enumerate(visible_plan):
        dot_class = "done" if index in completed else "active" if index == current else ""
        status = st.session_state.monitoring_status.get(step.get("day", ""), "Pending")
        html += (
            '<div class="timeline-row">'
            f'<div class="timeline-dot {dot_class}"></div>'
            '<div class="task-card">'
            f'<strong>{escape(str(step.get("day", "")))}: {escape(str(step.get("title", "")))}</strong><br>'
            f'<span>{status_label(status)}</span>'
            '</div></div>'
        )
    st.markdown(html, unsafe_allow_html=True)


def dashboard_page() -> None:
    render_hero(t("dashboard_title"), t("dashboard_intro"))
    render_dashboard_metrics()

    left, right = st.columns([1.15, 0.85])
    with left:
        render_current_task_card()
        render_current_step_panel()
    with right:
        st.subheader(t("quick_actions"))
        action_col1, action_col2 = st.columns(2)
        if action_col1.button(t("start_planning"), use_container_width=True):
            navigate_to_page("Planning")
        if action_col2.button(t("continue_monitoring"), use_container_width=True):
            navigate_to_page("Monitoring")
        action_col3, action_col4 = st.columns(2)
        if action_col3.button(t("open_ai_coach"), use_container_width=True):
            navigate_to_page("AI Coach")
        if action_col4.button(t("open_reports"), use_container_width=True):
            navigate_to_page("Reports")
        render_material_library()

    render_progress_timeline(limit=6)

    if st.session_state.tasks:
        st.subheader(t("recent_tasks"))
        recent = sorted(
            st.session_state.tasks.values(),
            key=lambda item: item.get("updated_at", ""),
            reverse=True,
        )[:4]
        for task in recent:
            title = escape(str(task.get("title", t("unsaved_task"))))
            updated = escape(str(task.get("updated_at", ""))[:16].replace("T", " "))
            st.markdown(
                f'<div class="task-card"><strong>{title}</strong><br><span>{t("updated_at")}: {updated}</span></div>',
                unsafe_allow_html=True,
            )


def task_library_page() -> None:
    render_hero(t("task_library_title"), t("task_library_intro"))
    task_ids = sorted(
        st.session_state.tasks.keys(),
        key=lambda item: st.session_state.tasks[item].get("updated_at", ""),
        reverse=True,
    )
    st.metric(t("saved_task_count"), len(task_ids))
    if not task_ids:
        st.info(t("no_saved_tasks"))
        if st.button(t("start_planning")):
            navigate_to_page("Planning")
        return

    for task_id in task_ids:
        task = st.session_state.tasks.get(task_id, {})
        title = escape(str(task.get("title", t("unsaved_task"))))
        created = escape(str(task.get("created_at", ""))[:16].replace("T", " "))
        updated = escape(str(task.get("updated_at", ""))[:16].replace("T", " "))
        current_badge = " · Current" if current_language() == "English" else " · 当前任务"
        is_current = task_id == st.session_state.get("current_task_id", "")
        st.markdown(
            f"""
            <div class="task-card">
                <strong>{title}{current_badge if is_current else ""}</strong><br>
                <span>{t("created_at")}: {created}</span><br>
                <span>{t("updated_at")}: {updated}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )

        with st.expander(t("task_preview"), expanded=False):
            st.text(task_preview_text(task_id))

        action_col1, action_col2, action_col3 = st.columns([1, 1, 1])
        if action_col1.button(t("load_task"), key=f"library_load_{task_id}", use_container_width=True):
            load_task_into_state(task_id)
            set_sidebar_flash("task_loaded")
            navigate_to_page("Dashboard")

        new_title = action_col2.text_input(
            t("task_name"),
            value=str(task.get("title", "")),
            key=f"library_rename_title_{task_id}",
        )
        if action_col2.button(t("rename_task"), key=f"library_rename_{task_id}", use_container_width=True):
            if rename_task(task_id, new_title):
                set_sidebar_flash("task_renamed")
                st.rerun()

        confirm_delete = action_col3.checkbox(
            t("confirm_delete_task"),
            key=f"library_confirm_delete_{task_id}",
        )
        if action_col3.button(
            t("delete_task"),
            key=f"library_delete_{task_id}",
            use_container_width=True,
            disabled=not confirm_delete,
        ):
            if delete_task(task_id):
                set_sidebar_flash("task_deleted")
                st.rerun()


def today_learning_page() -> None:
    render_hero(t("today_learning"), t("monitoring_intro"))
    if not st.session_state.study_plan:
        st.info(t("monitoring_need_plan"))
        if st.button(t("start_planning")):
            navigate_to_page("Planning")
        return
    render_workflow_progress()
    render_current_step_panel()
    st.subheader(t("quick_actions"))
    col1, col2, col3 = st.columns(3)
    if col1.button(t("continue_monitoring"), use_container_width=True):
        navigate_to_page("Monitoring")
    if col2.button(t("go_to_control"), use_container_width=True):
        navigate_to_page("Control")
    if col3.button(t("go_to_reflection"), use_container_width=True):
        navigate_to_page("Reflection")
    render_progress_timeline(limit=8)


def task_details_page() -> None:
    render_hero(t("task_details"), t("dashboard_intro"))
    render_current_task_card()
    if st.session_state.study_plan:
        render_dashboard_metrics()
        st.subheader(t("structured_plan"))
        render_plan(st.session_state.study_plan)
        strategies = st.session_state.plan_context.get("strategies", [])
        if strategies:
            st.subheader(t("recommended_strategies"))
            for strategy in strategies:
                st.markdown(f"- {strategy}")
    else:
        st.info(t("generate_plan_info"))
    render_material_upload_panel("task_details")


def reports_page() -> None:
    render_hero(t("reports_title"), t("reports_intro"))
    render_dashboard_metrics()
    if not st.session_state.workflow_step_reflections:
        st.info(t("no_reports_yet"))
    render_step_reflection_archive()
    render_progress_timeline()
    if st.session_state.learning_materials:
        render_material_library()


def render_numbered_guide(items: List[str]) -> None:
    for index, item in enumerate(items, start=1):
        st.markdown(
            f"""
            <div class="task-card">
                <small>{index}</small><br>
                <span>{escape(str(item))}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )


def guide_page() -> None:
    render_hero(t("guide_title"), t("guide_intro"))

    st.subheader(t("guide_quick_start"))
    render_numbered_guide(TRANSLATIONS["guide_quick_items"][current_language()])

    st.subheader(t("guide_srl_flow"))
    render_numbered_guide(TRANSLATIONS["guide_srl_items"][current_language()])

    demo_col, tip_col = st.columns(2)
    with demo_col:
        st.subheader(t("guide_class_demo"))
        render_numbered_guide(TRANSLATIONS["guide_demo_items"][current_language()])
    with tip_col:
        st.subheader(t("guide_tips"))
        render_numbered_guide(TRANSLATIONS["guide_tip_items"][current_language()])

    st.subheader(t("quick_actions"))
    col1, col2, col3 = st.columns(3)
    if col1.button(t("start_planning"), use_container_width=True):
        navigate_to_page("Planning")
    if col2.button(t("continue_monitoring"), use_container_width=True):
        navigate_to_page("Monitoring")
    if col3.button(t("open_ai_coach"), use_container_width=True):
        navigate_to_page("AI Coach")


def plan_length_range(deadline: date, difficulty: str) -> Tuple[int, int]:
    days_until_deadline = max((deadline - date.today()).days + 1, 1)
    if days_until_deadline <= 14:
        daily_items = max(1, min(days_until_deadline, 15))
        return daily_items, daily_items
    ranges = {"Low": (4, 8), "Medium": (5, 12), "High": (6, 15)}
    return ranges[difficulty]


def calculate_plan_length(deadline: date, daily_hours: int, difficulty: str) -> int:
    min_items, max_items = plan_length_range(deadline, difficulty)
    if min_items == max_items:
        return min_items

    baseline = {"Low": 6, "Medium": 8, "High": 10}[difficulty]
    weeks = weeks_until_deadline(deadline)
    if weeks >= 10:
        baseline += 1
    elif weeks <= 4:
        baseline -= 1

    if daily_hours <= 2:
        baseline += 1
    elif daily_hours >= 5:
        baseline -= 1

    return min(max(baseline, min_items), max_items, 15)


def plan_timeline_mode(deadline: date) -> str:
    days_until_deadline = max((deadline - date.today()).days + 1, 1)
    if days_until_deadline <= 14:
        return "daily"
    return "flexible"


def weeks_until_deadline(deadline: date) -> int:
    days_until_deadline = max((deadline - date.today()).days + 1, 1)
    return max(1, math.ceil(days_until_deadline / 7))


def plan_period_label(index: int, deadline: date) -> str:
    return date_range_label(index, index, deadline)


def format_plan_date(value: date) -> str:
    if current_language() == "中文":
        return f"{value.month}月{value.day}日"
    return f"{value.strftime('%b')} {value.day}"


def date_range_label(index: int, total_items: int, deadline: date) -> str:
    today = date.today()
    total_days = max((deadline - today).days + 1, 1)
    total_items = max(total_items, 1)
    start_offset = round((index - 1) * total_days / total_items)
    end_offset = round(index * total_days / total_items) - 1
    start_date = today + timedelta(days=max(start_offset, 0))
    end_date = today + timedelta(days=min(max(end_offset, start_offset), total_days - 1))

    if start_date == end_date:
        return format_plan_date(start_date)
    return f"{format_plan_date(start_date)} - {format_plan_date(end_date)}"


def plan_timeline_instruction(deadline: date) -> str:
    days_until_deadline = max((deadline - date.today()).days + 1, 1)
    mode = plan_timeline_mode(deadline)
    if mode == "daily":
        return (
            f"The deadline is {days_until_deadline} days away. Create a date-labeled daily plan that covers every day until the deadline."
        )
    return (
        f"The deadline is {days_until_deadline} days away. Create a plan using real calendar date-range labels in the day field, such as Apr 20 - Apr 26. The date ranges must cover the full period from today to the deadline, not only the first few days. Do not use Week 1/Week 2 or Phase 1/Phase 2 labels."
    )


def normalize_plan_periods(plan: List[Dict[str, str]], deadline: date) -> List[Dict[str, str]]:
    normalized = []
    total_items = max(len(plan), 1)
    for index, step in enumerate(plan, start=1):
        updated = dict(step)
        updated["day"] = date_range_label(index, total_items, deadline)
        normalized.append(updated)
    return normalized


def challenge_scaffold(challenge: str) -> Tuple[str, str]:
    challenge_text = challenge.lower()
    scaffolds = {
        "Too many materials": (
            "Use skimming, source triage, and synthesis tables before deep reading.",
            "Limit deep work to materials that directly support the current learning goal.",
        ),
        "No structure": (
            "Build a structure before producing the final output.",
            "Turn milestones into concrete actions before adding more detail.",
        ),
        "Task anxiety": (
            "Start with low-pressure first attempts and small output targets.",
            "Separate idea generation from editing to reduce perfectionism.",
        ),
        "Time pressure": (
            "Prioritise high-impact learning outputs first.",
            "Compress lower-priority reading and reserve time for revision.",
        ),
        "Not sure where to start": (
            "Begin with task interpretation and a simple question map.",
            "Define the goal, expected output, audience, and required evidence or performance standard.",
        ),
    }
    if challenge in scaffolds:
        return scaffolds[challenge]
    if any(word in challenge_text for word in ["reading", "readings", "source", "sources", "article", "articles"]):
        return scaffolds["Too many materials"]
    if any(word in challenge_text for word in ["structure", "outline", "argument", "organise", "organize"]):
        return scaffolds["No structure"]
    if any(word in challenge_text for word in ["anxiety", "anxious", "stress", "afraid", "confidence"]):
        return scaffolds["Task anxiety"]
    if any(word in challenge_text for word in ["time", "deadline", "late", "rush", "pressure"]):
        return scaffolds["Time pressure"]
    if any(word in challenge_text for word in ["start", "begin", "where", "unclear", "lost"]):
        return scaffolds["Not sure where to start"]
    return (
        "Use metacognitive checkpoints to diagnose the exact barrier before adding more work time.",
        "Convert the challenge into one observable next action and review it after the session.",
    )


def generate_study_plan(
    task_description: str,
    deadline: date,
    daily_hours: int,
    difficulty: str,
    challenge: str,
) -> Tuple[List[Dict[str, str]], List[str]]:
    plan_length = calculate_plan_length(deadline, daily_hours, difficulty)
    scaffold, priority_note = challenge_scaffold(challenge)

    base_steps = [
        (
            "Clarify the learning goal and success criteria",
            "Identify the target outcome, required deliverables, assessment criteria, and constraints.",
        ),
        (
            "Preview key materials and identify priority areas",
            "Skim readings, notes, examples, slides, or resources to decide what deserves focused attention.",
        ),
        (
            "Organise information into themes or skill areas",
            "Group concepts, evidence, examples, or practice problems so the task becomes easier to navigate.",
        ),
        (
            "Create a working structure or learning map",
            "Turn the goal into sections, milestones, concept links, practice targets, or output checkpoints.",
        ),
        (
            "Produce the first concrete output",
            "Create a draft section, solved example, concept summary, presentation outline, or practice response.",
        ),
        (
            "Complete the core learning tasks",
            "Work through the remaining sections, concepts, examples, or practice items while checking alignment with the goal.",
        ),
        (
            "Review understanding and improve weak areas",
            "Check accuracy, coherence, gaps, misconceptions, and areas that need another strategy.",
        ),
        (
            "Finalise, rehearse, or prepare for submission",
            "Polish the final output, test recall, rehearse performance, or prepare submission requirements.",
        ),
    ]

    if challenge in {"Writing anxiety", "Task anxiety"}:
        base_steps.insert(
            1,
            (
                "Complete a 20-minute confidence draft",
                "Create rough notes or a first attempt without editing to lower the barrier to starting.",
            ),
        )
    elif challenge == "No structure":
        base_steps.insert(
            2,
            (
                "Create a working structure and milestone hierarchy",
                "Write a provisional central idea and three to five milestones or sections.",
            ),
        )
    elif challenge in {"Too many readings", "Too many materials"}:
        base_steps.insert(
            2,
            (
                "Prioritise materials into core, useful, and optional groups",
                "Select the highest-value materials before committing time to deep work.",
            ),
        )
    elif challenge == "Time pressure":
        base_steps.insert(
            1,
            (
                "Define a minimum viable submission path",
                "Identify the learning outputs that matter most and the tasks that can be shortened.",
            ),
        )
    elif challenge == "Not sure where to start":
        base_steps.insert(
            1,
            (
                "Turn the task into three guiding questions",
                "Ask what the paper must explain, compare, evaluate, or propose.",
            ),
        )

    selected_steps = compress_steps(base_steps, plan_length)
    plan = []
    for index, (title, detail) in enumerate(selected_steps, start=1):
        plan.append(
            {
                "day": date_range_label(index, len(selected_steps), deadline),
                "title": title,
                "detail": detail,
                "time": f"{daily_hours} hour{'s' if daily_hours > 1 else ''}",
            }
        )

    strategies = [
        "Preview before deep work to protect time and attention.",
        "Use structure-first planning before producing the final output.",
        "Set retrieval checkpoints: explain or perform the key idea without looking at notes.",
        scaffold,
        priority_note,
    ]

    if difficulty == "High":
        strategies.append("Add short daily reflection notes to detect confusion early.")
    if daily_hours <= 2:
        strategies.append("Use small study blocks with one concrete output per session.")

    return plan, strategies


def generate_generic_fallback_study_plan(
    task_description: str,
    deadline: date,
    daily_hours: int,
    difficulty: str,
) -> Tuple[List[Dict[str, str]], List[str]]:
    return generate_study_plan(
        task_description=task_description,
        deadline=deadline,
        daily_hours=daily_hours,
        difficulty=difficulty,
        challenge="",
    )


def compress_steps(
    steps: List[Tuple[str, str]], target_length: int
) -> List[Tuple[str, str]]:
    target_length = max(1, target_length)
    if target_length >= len(steps):
        expanded = steps[:]
        if current_language() == "中文":
            expansion_steps = [
                ("主动回忆检查", "合上材料解释关键概念、方法或任务要求，记录仍然说不清楚的部分。"),
                ("应用练习或样例测试", "用一个练习题、案例、段落或小产出来检验当前理解是否能迁移。"),
                ("错误与障碍诊断", "整理错误、卡点和时间消耗来源，区分是概念问题、策略问题还是任务范围问题。"),
                ("策略调节小循环", "根据监控结果压缩、重排或拆分下一步任务，并设定一个可观察产出。"),
                ("最终整合与质量检查", "把学习证据整合成清晰产出，检查准确性、结构和提交或表现要求。"),
                ("下一轮学习准备", "保留有效策略，列出下一轮最小目标和需要优先澄清的问题。"),
            ]
        else:
            expansion_steps = [
                ("Active recall check", "Close the materials and explain the key concept, method, or requirement; note what remains unclear."),
                ("Application or example test", "Use one practice item, case, paragraph, or small output to test whether understanding transfers."),
                ("Error and barrier diagnosis", "List errors, stuck points, and time drains; separate concept issues, strategy issues, and scope issues."),
                ("Strategy adjustment mini-cycle", "Use monitoring evidence to compress, reorder, or break down the next task with one visible output."),
                ("Final integration and quality check", "Combine learning evidence into a clear output and check accuracy, structure, and requirements."),
                ("Prepare the next learning cycle", "Keep effective strategies and name the smallest next goal plus the highest-priority uncertainty."),
            ]
        while len(expanded) < target_length:
            expanded.append(expansion_steps[(len(expanded) - len(steps)) % len(expansion_steps)])
        return expanded

    compressed = steps[:]
    while len(compressed) > target_length:
        merge_index = 1 if len(compressed) > 3 else 0
        first = compressed.pop(merge_index)
        second = compressed.pop(merge_index)
        if current_language() == "中文":
            title = f"{first[0]} + {second[0]}"
        else:
            title = f"{first[0]} and {second[0].lower()}"
        compressed.insert(merge_index, (title, f"{first[1]} {second[1]}"))
    return compressed


def build_openai_client():
    base_url = get_openai_base_url()
    client_kwargs = {"api_key": get_openai_api_key()}
    if base_url:
        client_kwargs["base_url"] = base_url
    return OpenAI(**client_kwargs)


def call_genai_text(system_prompt: str, user_prompt: str, temperature: float = 0.4) -> str:
    if not has_openai_key():
        raise RuntimeError("OpenAI API key is not configured.")

    client = build_openai_client()
    response = client.chat.completions.create(
        model=get_openai_model(),
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=temperature,
    )
    return response.choices[0].message.content or ""


def call_genai_json(
    system_prompt: str, user_prompt: str, temperature: float = 0.35
) -> Dict[str, object]:
    if not has_openai_key():
        raise RuntimeError("OpenAI API key is not configured.")

    client = build_openai_client()
    response = client.chat.completions.create(
        model=get_openai_model(),
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=temperature,
        response_format={"type": "json_object"},
    )
    content = response.choices[0].message.content or "{}"
    return json.loads(content)


def source_caption(source: str) -> None:
    if source:
        st.caption(f"{t('output_source')}: {translated_source(source)}")


def generate_study_plan_with_genai(
    task_description: str,
    deadline: date,
    daily_hours: int,
    difficulty: str,
    challenge: str,
    optional_context: Dict[str, str] | None = None,
) -> Tuple[List[Dict[str, str]], List[str], str, str]:
    optional_context = optional_context or {}
    fallback_plan, fallback_strategies = generate_generic_fallback_study_plan(
        task_description, deadline, daily_hours, difficulty
    )
    target_days = calculate_plan_length(deadline, daily_hours, difficulty)
    min_plan_items, max_plan_items = plan_length_range(deadline, difficulty)
    timeline_instruction = plan_timeline_instruction(deadline)

    system_prompt = (
        "You are PlanWise, an LLM-supported Self-Regulated Learning planner for "
        "university learning, study goals, and academic tasks. Generate practical study plans grounded in "
        f"Zimmerman's SRL framework. Write all user-facing JSON values in {output_language_name()}. "
        "Return valid JSON only."
    )
    user_prompt = f"""
Create a tailored study plan for a learning challenge, study goal, or academic task.

Inputs:
- Task description: {task_description}
- Deadline: {deadline}
- Daily available study time: {daily_hours} hours
- Perceived difficulty: {difficulty}
- Student-described current challenge: {challenge or "Not specified"}
- Optional learner/task context: {json.dumps(optional_context, ensure_ascii=False)}
- Suggested fallback number of plan items: {target_days}
- Acceptable plan item range: {min_plan_items} to {max_plan_items}
- Timeline instruction: {timeline_instruction}

Return JSON with this exact structure:
{{
  "plan": [
    {{
      "day": "Apr 20 - Apr 26",
      "title": "short action title",
      "detail": "specific learning action with SRL/metacognitive support",
      "time": "{daily_hours} hours"
    }}
  ],
  "strategies": [
    "strategy 1",
    "strategy 2",
    "strategy 3"
  ],
  "srl_rationale": "one concise sentence explaining how the plan supports Planning, Monitoring, Control, and Reflection"
}}

Requirements:
- Write all user-facing JSON values in {output_language_name()}.
- Choose a suitable number of plan items within the acceptable range: {min_plan_items} to {max_plan_items}.
- Prefer approximately {target_days} plan items. Do not use the upper bound just because the perceived difficulty is high.
- Use fewer items for simpler or more focused goals, and more items for complex, high-stakes, or multi-material goals.
- Avoid repetitive checkpoint-only items. Each item should move the learner forward with a distinct goal, activity, or SRL strategy.
- The plan must cover the full time from today to the deadline, not only the first week.
- Use real calendar time labels in the "day" field, preferably date ranges such as "Apr 20 - Apr 26".
- Do not use Week 1/Week 2 or Phase 1/Phase 2 labels.
- Keep each day concrete and student-facing.
- Use the optional context only when it is provided; do not invent missing requirements.
- Include learning strategies such as previewing, active recall, concept mapping, spaced practice, synthesis, output-first work, revision, or reflection.
- Respond to the student's own challenge wording instead of forcing it into a preset category.
- Do not include Markdown fences.
"""
    try:
        data = call_genai_json(system_prompt, user_prompt)
        plan = data.get("plan", [])
        strategies = data.get("strategies", [])
        rationale = str(data.get("srl_rationale", "")).strip()
        if not isinstance(plan, list) or not plan:
            raise ValueError("GenAI response did not include a usable plan.")
        if not isinstance(strategies, list) or not strategies:
            strategies = fallback_strategies

        cleaned_plan = []
        for index, step in enumerate(plan[:max_plan_items], start=1):
            if not isinstance(step, dict):
                continue
            cleaned_plan.append(
                {
                    "day": str(step.get("day") or f"Day {index}"),
                    "title": str(step.get("title") or fallback_plan[min(index - 1, len(fallback_plan) - 1)]["title"]),
                    "detail": str(step.get("detail") or fallback_plan[min(index - 1, len(fallback_plan) - 1)]["detail"]),
                    "time": str(step.get("time") or f"{daily_hours} hours"),
                }
            )
        if not cleaned_plan:
            raise ValueError("GenAI response plan was empty after validation.")
        while len(cleaned_plan) < min_plan_items and len(cleaned_plan) < len(fallback_plan):
            cleaned_plan.append(fallback_plan[len(cleaned_plan)])

        cleaned_plan = normalize_plan_periods(cleaned_plan, deadline)
        return cleaned_plan, [str(item) for item in strategies], "GenAI-generated via OpenAI API", rationale
    except Exception as exc:
        return (
            fallback_plan,
            fallback_strategies,
            "Rule-based fallback after GenAI was unavailable",
            f"Fallback reason: {exc}",
        )


def render_plan(plan: List[Dict[str, str]]) -> None:
    for step in plan:
        day = escape(str(step["day"]))
        title = escape(str(step["title"]))
        detail = escape(str(step["detail"]))
        time = escape(str(step["time"]))
        st.markdown(
            f"""
            <div class="plan-day">
                <strong>{day}: {title}</strong><br>
                <span>{detail}</span><br>
                <small>{'建议学习时间' if current_language() == '中文' else 'Suggested study time'}: {time}</small>
            </div>
            """,
            unsafe_allow_html=True,
        )


def planning_page() -> None:
    render_header("Planning")
    st.markdown(f'<div class="phase-label">{t("phase_prefix")}: {phase_label("Planning")}</div>', unsafe_allow_html=True)
    st.subheader(t("goal_setting"))
    st.write(t("planning_intro"))
    render_material_upload_panel("planning")

    with st.form("planning_form"):
        task_description = st.text_area(
            t("learning_goal_or_task"),
            placeholder=t("learning_goal_placeholder"),
            height=110,
        )
        col1, col2 = st.columns(2)
        with col1:
            deadline = st.date_input(t("deadline"), min_value=date.today())
            difficulty = st.selectbox(
                t("perceived_difficulty"),
                ["Low", "Medium", "High"],
                index=1,
                format_func=difficulty_label,
            )
        with col2:
            daily_hours = st.slider(t("daily_time"), 1, 6, 3)
            challenge = st.text_area(
                t("current_challenge"),
                placeholder=t("challenge_placeholder"),
                height=95,
            )

        with st.expander(t("optional_learning_context"), expanded=True):
            st.caption(t("optional_context_caption"))
            opt_col1, opt_col2 = st.columns(2)
            with opt_col1:
                assignment_type = st.text_input(
                    t("task_type"),
                    placeholder=t("task_type_placeholder"),
                )
                word_count = st.text_input(
                    t("expected_output"),
                    placeholder=t("expected_output_placeholder"),
                )
                discipline = st.text_input(
                    t("course_or_discipline"),
                    placeholder=t("course_placeholder"),
                )
                required_sources = st.text_input(
                    t("required_materials"),
                    placeholder=t("required_materials_placeholder"),
                )
            with opt_col2:
                grading_criteria = st.text_area(
                    t("success_criteria"),
                    placeholder=t("success_criteria_placeholder"),
                    height=90,
                )
                current_progress = st.text_area(
                    t("current_progress"),
                    placeholder=t("current_progress_placeholder"),
                    height=90,
                )
                existing_materials = st.text_area(
                    t("existing_materials"),
                    placeholder=t("existing_materials_placeholder"),
                    height=90,
                )
        submitted = st.form_submit_button(t("generate_study_plan"), type="primary")

    if submitted:
        if not task_description.strip():
            st.warning(t("task_required_warning"))
        else:
            optional_context = {
                "assignment_type": assignment_type.strip(),
                "word_count": word_count.strip(),
                "discipline": discipline.strip(),
                "required_sources": required_sources.strip(),
                "grading_criteria": grading_criteria.strip(),
                "current_progress": current_progress.strip(),
                "existing_materials": existing_materials.strip(),
            }
            optional_context = {
                key: value for key, value in optional_context.items() if value
            }
            material_brief = materials_brief()
            if material_brief:
                optional_context["uploaded_learning_materials"] = material_brief
            with st.spinner(t("planning_spinner")):
                plan, strategies, source, rationale = generate_study_plan_with_genai(
                    task_description.strip(),
                    deadline,
                    daily_hours,
                    difficulty,
                    challenge.strip(),
                    optional_context,
                )
            st.session_state.study_plan = plan
            st.session_state.plan_context = {
                "task_description": task_description.strip(),
                "deadline": str(deadline),
                "daily_hours": daily_hours,
                "difficulty": difficulty,
                "challenge": challenge.strip(),
                "optional_context": optional_context,
                "materials_brief": material_brief,
                "strategies": strategies,
                "planning_source": source,
                "planning_rationale": rationale,
            }
            st.session_state.monitoring_status = {
                step["day"]: "Pending" for step in plan
            }
            st.session_state.monitoring_details = {
                step["day"]: {
                    "evidence": "",
                    "barrier": "",
                    "confidence": 3,
                    "actual_time": float(daily_hours),
                }
                for step in plan
            }
            for key in list(st.session_state.keys()):
                if (
                    str(key).startswith("check_")
                    or str(key).startswith("evidence_")
                    or str(key).startswith("barrier_")
                    or str(key).startswith("confidence_")
                    or str(key).startswith("actual_time_")
                ):
                    del st.session_state[key]
            st.session_state.adjusted_plan = []
            st.session_state.monitoring_feedback = {}
            st.session_state.control_recommendations = []
            st.session_state.control_next_actions = []
            st.session_state.control_context = {}
            st.session_state.control_source = ""
            st.session_state.reflection_report = {}
            st.session_state.reflection_inputs = {}
            st.session_state.workflow_current_index = 0
            st.session_state.workflow_completed_indices = []
            st.session_state.workflow_step_reflections = {}
            save_current_task(create_if_missing=True)
            st.success(t("plan_success"))

    if st.session_state.study_plan:
        st.subheader(t("structured_plan"))
        source_caption(st.session_state.plan_context.get("planning_source", ""))
        render_plan(st.session_state.study_plan)
        rationale = st.session_state.plan_context.get("planning_rationale", "")
        if rationale:
            st.info(rationale)
        st.subheader(t("recommended_strategies"))
        for strategy in st.session_state.plan_context.get("strategies", []):
            st.markdown(f"- {strategy}")
    else:
        st.info(t("generate_plan_info"))


def workflow_steps_through_current() -> List[Dict[str, str]]:
    plan = st.session_state.get("study_plan", [])
    if not plan:
        return []
    return plan[: current_step_index() + 1]


def workflow_step_days_through_current() -> set[str]:
    return {str(step.get("day", "")) for step in workflow_steps_through_current()}


def monitoring_status_through_current() -> Dict[str, str]:
    active_days = workflow_step_days_through_current()
    return {
        day: status
        for day, status in st.session_state.monitoring_status.items()
        if day in active_days
    }


def monitoring_details_through_current() -> Dict[str, Dict[str, object]]:
    active_days = workflow_step_days_through_current()
    return {
        day: details
        for day, details in st.session_state.monitoring_details.items()
        if day in active_days
    }


def get_status_counts(through_current: bool = False) -> Dict[str, int]:
    counts = {status: 0 for status in STATUS_OPTIONS}
    statuses = (
        monitoring_status_through_current()
        if through_current
        else st.session_state.monitoring_status
    )
    for status in statuses.values():
        counts[status] = counts.get(status, 0) + 1
    return counts


def monitoring_completion_rate(counts: Dict[str, int]) -> float:
    total_tracked = sum(counts.values()) - counts.get("Pending", 0)
    if total_tracked <= 0:
        return 0.0
    return round((counts.get("Completed", 0) / total_tracked) * 100, 1)


def adaptive_feedback(counts: Dict[str, int], challenge: str = "") -> List[str]:
    feedback = []
    if counts.get("Delayed", 0) >= 1:
        feedback.append(
            "你可能需要减少被动复习，优先完成一个看得见的学习产出。"
            if current_language() == "中文"
            else "You may need to reduce passive review and prioritise a visible learning output."
        )
    if counts.get("Confused", 0) >= 1:
        feedback.append(
            "可以考虑改用综合型笔记和概念图来整理理解。"
            if current_language() == "中文"
            else "Consider switching to synthesis-based note taking and concept mapping."
        )
    if counts.get("Overwhelmed", 0) >= 1:
        feedback.append(
            "请减少材料负荷，并为下一次学习设定一个具体产出。"
            if current_language() == "中文"
            else "Reduce the material load and define one concrete output for the next session."
        )
    if counts.get("Completed", 0) >= 2 and counts.get("Delayed", 0) == 0:
        feedback.append(
            "你的监控信号比较积极。继续使用回忆检查点来测试理解。"
            if current_language() == "中文"
            else "Your monitoring signals are positive. Keep using retrieval checkpoints to test understanding."
        )
    if challenge == "No structure" and counts.get("Confused", 0) >= 1:
        feedback.append(
            "在加入更多材料前，先回到目标、结构和关键检查点。"
            if current_language() == "中文"
            else "Return to the goal, structure, and key checkpoints before adding more materials."
        )
    if not feedback:
        feedback.append(
            "至少记录一个学习步骤后，PlanWise 才能提供适应性自我调节学习反馈。"
            if current_language() == "中文"
            else "Track at least one study step to receive adaptive SRL feedback."
        )
    return feedback


def monitoring_feedback_with_genai(
    counts: Dict[str, int],
    challenge: str,
    plan: List[Dict[str, str]],
    statuses: Dict[str, str],
    details: Dict[str, Dict[str, object]],
) -> Tuple[List[str], str]:
    fallback = adaptive_feedback(counts, challenge)
    system_prompt = (
        "You are PlanWise, an LLM-supported SRL monitoring coach. Interpret student "
        "progress signals for learning tasks and return concise adaptive feedback. "
        f"Write all user-facing JSON values in {output_language_name()}. Return valid JSON only."
    )
    user_prompt = f"""
The learner is monitoring a study plan.

Plan context:
{json.dumps(st.session_state.plan_context, ensure_ascii=False)}

Uploaded learning materials brief:
{materials_brief()}

Study plan:
{json.dumps(plan, ensure_ascii=False)}

Monitoring statuses by step:
{json.dumps(statuses, ensure_ascii=False)}

Monitoring details by step:
{json.dumps(details, ensure_ascii=False)}

Status counts:
{json.dumps(counts, ensure_ascii=False)}

Return JSON with this exact structure:
{{
  "feedback": [
    "specific monitoring feedback item",
    "specific monitoring feedback item"
  ]
}}

Requirements:
- Write all feedback in {output_language_name()}.
- Feedback must reference monitoring awareness, time use, confusion, delay, or overload.
- Use the learner's evidence, barrier notes, confidence ratings, and actual time when available.
- Give 2 to 4 concise items.
- Keep the voice supportive and educational.
- Do not include Markdown fences.
"""
    try:
        data = call_genai_json(system_prompt, user_prompt)
        feedback = data.get("feedback", [])
        if not isinstance(feedback, list) or not feedback:
            raise ValueError("GenAI response did not include feedback.")
        return [str(item) for item in feedback[:4]], "GenAI-generated via OpenAI API"
    except Exception as exc:
        return (
            fallback + [f"Fallback note: GenAI monitoring feedback was unavailable ({exc})."],
            "Rule-based fallback after GenAI was unavailable",
        )


def render_monitoring_step_editor(step: Dict[str, str], expanded: bool = True, key_prefix: str = "") -> None:
    current = st.session_state.monitoring_status.get(step["day"], "Pending")
    step_details = st.session_state.monitoring_details.get(
        step["day"],
        {
            "evidence": "",
            "barrier": "",
            "confidence": 3,
            "actual_time": float(st.session_state.plan_context.get("daily_hours", 1)),
        },
    )
    prefix = f"{key_prefix}_{step['day']}" if key_prefix else step["day"]
    with st.expander(f"{step['day']}: {step['title']}", expanded=expanded):
        st.caption(step["detail"])
        cols = st.columns(4)
        selected_signals = []
        for index, status in enumerate(["Completed", "Delayed", "Confused", "Overwhelmed"]):
            key = f"check_{prefix}_{status}"
            if key not in st.session_state:
                st.session_state[key] = current == status
            checked = cols[index].checkbox(status_label(status), key=key)
            if checked:
                selected_signals.append(status)

        if selected_signals:
            priority = ["Overwhelmed", "Confused", "Delayed", "Completed"]
            primary_status = next(status for status in priority if status in selected_signals)
            st.session_state.monitoring_status[step["day"]] = primary_status
            if len(selected_signals) > 1:
                st.caption(t("multi_signal_caption").format(status=status_label(primary_status)))
        else:
            st.session_state.monitoring_status[step["day"]] = "Pending"

        detail_col1, detail_col2 = st.columns(2)
        with detail_col1:
            evidence = st.text_area(
                t("evidence_progress"),
                value=str(step_details.get("evidence", "")),
                placeholder=t("evidence_placeholder"),
                height=85,
                key=f"evidence_{prefix}",
            )
            confidence = st.slider(
                t("confidence_after_step"),
                1,
                5,
                int(step_details.get("confidence", 3)),
                key=f"confidence_{prefix}",
            )
        with detail_col2:
            barrier = st.text_area(
                t("barrier_point"),
                value=str(step_details.get("barrier", "")),
                placeholder=t("barrier_placeholder"),
                height=85,
                key=f"barrier_{prefix}",
            )
            actual_time = st.number_input(
                t("actual_time"),
                min_value=0.0,
                max_value=12.0,
                value=float(step_details.get("actual_time", st.session_state.plan_context.get("daily_hours", 1))),
                step=0.5,
                key=f"actual_time_{prefix}",
            )

        st.session_state.monitoring_details[step["day"]] = {
            "evidence": evidence.strip(),
            "barrier": barrier.strip(),
            "confidence": confidence,
            "actual_time": actual_time,
        }


def monitoring_page() -> None:
    render_header("Monitoring")
    st.markdown(f'<div class="phase-label">{t("phase_prefix")}: {phase_label("Monitoring")}</div>', unsafe_allow_html=True)
    st.subheader(t("monitoring_dashboard"))

    if not st.session_state.study_plan:
        st.info(t("monitoring_need_plan"))
        return

    render_workflow_progress()
    st.write(t("monitoring_intro"))

    plan = st.session_state.study_plan
    current_index = current_step_index()
    step = current_step()
    if step:
        st.subheader(f"{t('current_step')}: {step['day']}")
        render_monitoring_step_editor(step, expanded=True, key_prefix="current")

    if current_index > 0:
        with st.expander(t("previous_steps"), expanded=False):
            st.caption(t("read_only_note"))
            edit_previous = st.checkbox(t("edit_previous_steps"), value=False)
            for index, previous_step in enumerate(plan[:current_index]):
                if edit_previous:
                    render_monitoring_step_editor(previous_step, expanded=False, key_prefix=f"previous_{index}")
                else:
                    render_step_summary(previous_step, index)

    if current_index < len(plan) - 1:
        with st.expander(t("upcoming_steps"), expanded=False):
            for index, future_step in enumerate(plan[current_index + 1 :], start=current_index + 1):
                render_step_summary(future_step, index)

    counts = get_status_counts(through_current=True)
    col1, col2, col3, col4, col5 = st.columns(5)
    col1.metric(t("completed"), counts.get("Completed", 0))
    col2.metric(t("delayed"), counts.get("Delayed", 0))
    col3.metric(t("confused"), counts.get("Confused", 0))
    col4.metric(t("overwhelmed"), counts.get("Overwhelmed", 0))
    col5.metric(t("tracked_completion"), f"{monitoring_completion_rate(counts)}%")

    st.subheader(t("adaptive_feedback"))
    challenge = st.session_state.plan_context.get("challenge", "")
    if st.button(t("generate_monitoring_feedback"), type="primary"):
        with st.spinner(t("monitoring_spinner")):
            feedback, source = monitoring_feedback_with_genai(
                counts,
                challenge,
                workflow_steps_through_current(),
                monitoring_status_through_current(),
                monitoring_details_through_current(),
            )
            st.session_state.monitoring_feedback = {
                "items": feedback,
                "source": source,
            }

    feedback_state = st.session_state.monitoring_feedback
    if not feedback_state:
        feedback, source = adaptive_feedback(counts, challenge), "Rule-based preview before GenAI generation"
    else:
        feedback, source = feedback_state.get("items", []), feedback_state.get("source", "")
    source_caption(source)
    for item in feedback:
        st.markdown(f'<div class="feedback-box">{escape(str(item))}</div>', unsafe_allow_html=True)

    st.subheader(t("monitoring_next_step_title"))
    if is_current_step_recorded():
        st.info(t("monitoring_next_step_hint"))
        nav_col1, nav_col2 = st.columns(2)
        if nav_col1.button(t("go_to_control"), use_container_width=True):
            save_current_task()
            navigate_to_page("Control")
        if nav_col2.button(t("go_to_reflection"), use_container_width=True):
            save_current_task()
            navigate_to_page("Reflection")
    else:
        st.caption(t("record_current_step_first"))
    save_current_task()


def generate_adjusted_plan() -> List[Dict[str, str]]:
    original_plan = st.session_state.study_plan
    statuses = st.session_state.monitoring_status
    adjusted = []

    for step in original_plan:
        status = statuses.get(step["day"], "Pending")
        title = step["title"]
        detail = step["detail"]

        if status == "Completed":
            detail = (
                "已完成。把这个产出作为下一步学习的基础。"
                if current_language() == "中文"
                else "Completed. Use this output as the foundation for the next step."
            )
        elif status == "Delayed":
            title = f"{'压缩版' if current_language() == '中文' else 'Compressed'}: {title}"
            detail = (
                "将这个任务缩短到最高价值的产出。跳过可选材料，只捕捉必要证据或例子，并尽快形成一个具体产出。"
                if current_language() == "中文"
                else "Shorten this task to its highest-value output. Skip optional materials, "
                "capture only essential evidence or examples, and move quickly toward a concrete output."
            )
        elif status == "Confused":
            title = f"{'先澄清再继续' if current_language() == '中文' else 'Clarify before continuing'}: {title}"
            detail = (
                "加入概念图步骤。在继续增加任务前，先明确关键概念、关系和不确定点。"
                if current_language() == "中文"
                else "Insert a concept mapping step. Define key concepts, relationships, "
                "and uncertainties before adding more work."
            )
        elif status == "Overwhelmed":
            title = f"{'减负版' if current_language() == '中文' else 'Reduced-load version'}: {title}"
            detail = (
                "缩小本次学习范围。选择一个资源、一个部分或一道练习，先完成这个小单元，再继续扩展。"
                if current_language() == "中文"
                else "Reduce scope for this session. Choose one resource, one section, or one practice item "
                "and finish that smaller unit before expanding."
            )

        adjusted.append(
            {
                "day": step["day"],
                "title": title,
                "detail": detail,
                "time": step["time"],
                "status": status,
            }
        )

    return adjusted


def control_recommendations(counts: Dict[str, int]) -> List[str]:
    if current_language() == "中文":
        recommendations = []
        if counts.get("Delayed", 0):
            recommendations.append("压缩低优先级任务，把时间保护给最关键的学习产出。")
        if counts.get("Confused", 0):
            recommendations.append("在下一次专注学习前加入概念图或过程图。")
        if counts.get("Overwhelmed", 0):
            recommendations.append("减少材料负荷，并定义一个最小可完成目标。")
        if counts.get("Pending", 0) and not any(counts.get(s, 0) for s in ["Delayed", "Confused", "Overwhelmed"]):
            recommendations.append("从最小的未记录步骤开始，并设定一个可见产出。")
        if not recommendations:
            recommendations.append("保持当前计划，并在完成后安排一次短反思。")
        return recommendations

    recommendations = []
    if counts.get("Delayed", 0):
        recommendations.append("Compress low-priority tasks and protect time for the most important learning output.")
    if counts.get("Confused", 0):
        recommendations.append("Insert a concept map before the next focused work session.")
    if counts.get("Overwhelmed", 0):
        recommendations.append("Reduce reading load and define a minimum viable section target.")
    if counts.get("Pending", 0) and not any(counts.get(s, 0) for s in ["Delayed", "Confused", "Overwhelmed"]):
        recommendations.append("Begin with the smallest pending step and set a visible output target.")
    if not recommendations:
        recommendations.append("Maintain the current plan and schedule a short reflection after submission.")
    return recommendations


def control_risk_snapshot(counts: Dict[str, int]) -> Dict[str, str]:
    challenge_signals = counts.get("Delayed", 0) + counts.get("Confused", 0) + counts.get("Overwhelmed", 0)
    if counts.get("Overwhelmed", 0) or challenge_signals >= 3:
        risk = t("high_risk")
    elif challenge_signals >= 1:
        risk = t("medium_risk")
    else:
        risk = t("low_risk")

    signal_order = ["Overwhelmed", "Confused", "Delayed", "Pending", "Completed"]
    primary = max(signal_order, key=lambda status: counts.get(status, 0))
    return {"risk": risk, "primary_signal": status_label(primary)}


def control_adjustment_with_genai(
    counts: Dict[str, int],
    plan: List[Dict[str, str]],
    statuses: Dict[str, str],
    control_context: Dict[str, object],
) -> Tuple[List[Dict[str, str]], List[str], List[str], str]:
    fallback_plan = generate_adjusted_plan()
    fallback_recommendations = control_recommendations(counts)
    fallback_next_actions = (
        [
            "选择一个最小可完成任务，并在 25 分钟内完成一个可见产出。",
            "学习后记录一个证据和一个仍然不清楚的问题。",
        ]
        if current_language() == "中文"
        else [
            "Choose one minimum viable task and complete one visible output within 25 minutes.",
            "After studying, record one piece of evidence and one remaining uncertainty.",
        ]
    )
    system_prompt = (
        "You are PlanWise, an LLM-supported SRL control-phase coach. Adjust academic "
        "study plans based on monitoring data. Return valid JSON only."
    )
    user_prompt = f"""
Revise the learner's study plan using SRL Control-phase logic.

Plan context:
{json.dumps(st.session_state.plan_context, ensure_ascii=False)}

Original study plan:
{json.dumps(plan, ensure_ascii=False)}

Current workflow step:
{json.dumps(current_step(), ensure_ascii=False)}

Monitoring statuses:
{json.dumps(statuses, ensure_ascii=False)}

Monitoring details:
{json.dumps(monitoring_details_through_current(), ensure_ascii=False)}

Learner adjustment preferences:
{json.dumps(control_context, ensure_ascii=False)}

Status counts:
{json.dumps(counts, ensure_ascii=False)}

Return JSON with this exact structure:
{{
  "adjusted_plan": [
    {{
      "day": "Day 1",
      "title": "adjusted action title",
      "detail": "specific adjustment to task, scope, strategy, or sequence",
      "time": "same or revised time",
      "status": "Completed/Delayed/Confused/Overwhelmed/Pending"
    }}
  ],
  "recommendations": [
    "control-phase recommendation"
  ],
  "next_actions": [
    "immediate action protocol for the next study session"
  ]
}}

Rules:
- Write all user-facing JSON values in {output_language_name()}.
- Use the learner adjustment preferences when they are provided.
- If delayed, compress low-priority tasks and protect output time.
- If confused, insert concept mapping, process mapping, or synthesis support.
- If overwhelmed, reduce reading load and shrink the next action.
- If completed, preserve momentum and connect outputs to the next task.
- Keep the adjusted plan the same length as the original plan.
- Provide 2 to 4 next_actions that are immediately executable.
- Do not include Markdown fences.
"""
    try:
        data = call_genai_json(system_prompt, user_prompt)
        adjusted_plan = data.get("adjusted_plan", [])
        recommendations = data.get("recommendations", [])
        next_actions = data.get("next_actions", [])
        if not isinstance(adjusted_plan, list) or not adjusted_plan:
            raise ValueError("GenAI response did not include an adjusted plan.")
        if not isinstance(recommendations, list) or not recommendations:
            recommendations = fallback_recommendations
        if not isinstance(next_actions, list) or not next_actions:
            next_actions = fallback_next_actions

        cleaned_plan = []
        for index, step in enumerate(adjusted_plan[: len(plan)]):
            if not isinstance(step, dict):
                continue
            original = plan[index]
            cleaned_plan.append(
                {
                    "day": str(step.get("day") or original.get("day", f"Day {index + 1}")),
                    "title": str(step.get("title") or original.get("title", "")),
                    "detail": str(step.get("detail") or original.get("detail", "")),
                    "time": str(step.get("time") or original.get("time", "")),
                    "status": str(step.get("status") or statuses.get(original.get("day", ""), "Pending")),
                }
            )
        if not cleaned_plan:
            raise ValueError("GenAI adjusted plan was empty after validation.")
        while len(cleaned_plan) < len(fallback_plan):
            cleaned_plan.append(fallback_plan[len(cleaned_plan)])

        return (
            cleaned_plan,
            [str(item) for item in recommendations[:5]],
            [str(item) for item in next_actions[:4]],
            "GenAI-generated via OpenAI API",
        )
    except Exception as exc:
        return (
            fallback_plan,
            fallback_recommendations + [f"Fallback note: GenAI control adjustment was unavailable ({exc})."],
            fallback_next_actions,
            "Rule-based fallback after GenAI was unavailable",
        )


def control_page() -> None:
    render_header("Control")
    st.markdown(f'<div class="phase-label">{t("phase_prefix")}: {phase_label("Control")}</div>', unsafe_allow_html=True)
    st.subheader(t("control_title"))

    if not st.session_state.study_plan:
        st.info(t("control_need_plan"))
        return

    st.write(t("control_intro"))
    render_workflow_progress()
    active_step = current_step()
    if active_step:
        st.subheader(f"{t('current_step')}: {active_step['day']}")
        render_step_summary(active_step, current_step_index())
    render_step_reflection_archive()

    counts = get_status_counts(through_current=True)
    snapshot = control_risk_snapshot(counts)
    snap_col1, snap_col2, snap_col3 = st.columns(3)
    snap_col1.metric(t("risk_level"), snapshot["risk"])
    snap_col2.metric(t("primary_signal"), snapshot["primary_signal"])
    snap_col3.metric(t("tracked_completion"), f"{monitoring_completion_rate(counts)}%")

    with st.expander(t("control_preferences"), expanded=True):
        st.caption(t("control_preferences_caption"))
        pref_col1, pref_col2 = st.columns(2)
        goal_options = TRANSLATIONS["control_goal_options"][current_language()]
        style_options = TRANSLATIONS["control_style_options"][current_language()]
        with pref_col1:
            adjustment_goal = st.selectbox(t("control_goal"), goal_options)
            adjustment_style = st.selectbox(t("control_style"), style_options)
        with pref_col2:
            next_session_time = st.slider(
                t("available_time_next"),
                0.5,
                6.0,
                float(st.session_state.plan_context.get("daily_hours", 1)),
                step=0.5,
            )
            non_negotiables = st.text_area(
                t("non_negotiables"),
                placeholder=t("non_negotiables_placeholder"),
                height=85,
            )

    control_context = {
        "adjustment_goal": adjustment_goal,
        "adjustment_style": adjustment_style,
        "available_time_for_next_session": next_session_time,
        "non_negotiables": non_negotiables.strip(),
        "risk_snapshot": snapshot,
    }
    st.session_state.control_context = control_context

    if st.button(t("generate_control_adjustment"), type="primary"):
        with st.spinner(t("control_spinner")):
            adjusted_plan, recommendations, next_actions, source = control_adjustment_with_genai(
                counts,
                st.session_state.study_plan,
                monitoring_status_through_current(),
                control_context,
            )
            st.session_state.adjusted_plan = adjusted_plan
            st.session_state.control_recommendations = recommendations
            st.session_state.control_next_actions = next_actions
            st.session_state.control_source = source

    st.subheader(t("updated_recommendations"))
    if st.session_state.control_recommendations:
        recommendations = st.session_state.control_recommendations
        source = st.session_state.control_source
    else:
        recommendations = control_recommendations(counts)
        source = "Rule-based preview before GenAI generation"
    source_caption(source)
    for item in recommendations:
        st.markdown(f'<div class="warning-box">{escape(str(item))}</div>', unsafe_allow_html=True)

    next_actions = st.session_state.control_next_actions
    if not next_actions:
        next_actions = (
            [
                "选择一个最小可完成任务，并在 25 分钟内完成一个可见产出。",
                "学习后记录一个证据和一个仍然不清楚的问题。",
            ]
            if current_language() == "中文"
            else [
                "Choose one minimum viable task and complete one visible output within 25 minutes.",
                "After studying, record one piece of evidence and one remaining uncertainty.",
            ]
        )
    st.subheader(t("next_actions"))
    for action in next_actions:
        st.markdown(f'<div class="feedback-box">{escape(str(action))}</div>', unsafe_allow_html=True)

    plan_to_show = st.session_state.adjusted_plan or generate_adjusted_plan()
    st.subheader(t("adjusted_plan"))
    if plan_to_show:
        if len(plan_to_show) == len(st.session_state.study_plan):
            display_steps = [plan_to_show[current_step_index()]]
        else:
            display_steps = plan_to_show
    else:
        display_steps = []
    for step in display_steps:
        day = escape(str(step["day"]))
        title = escape(str(step["title"]))
        detail = escape(str(step["detail"]))
        status = escape(str(step.get("status", "Pending")))
        st.markdown(
            f"""
            <div class="plan-day">
                <strong>{day}: {title}</strong><br>
                <span>{detail}</span><br>
                <small>{t("monitoring_status")}: {status_label(str(step.get("status", "Pending")))}</small>
            </div>
            """,
            unsafe_allow_html=True,
        )

    st.subheader(t("control_next_step_title"))
    st.info(t("control_next_step_hint"))
    nav_col1, nav_col2 = st.columns(2)
    if nav_col1.button(t("back_to_monitoring"), use_container_width=True):
        save_current_task()
        navigate_to_page("Monitoring")
    if nav_col2.button(t("go_to_reflection"), type="primary", use_container_width=True):
        save_current_task()
        navigate_to_page("Reflection")
    save_current_task()


def learner_profile(counts: Dict[str, int]) -> Dict[str, str]:
    total = max(sum(counts.values()), 1)
    completed_ratio = counts.get("Completed", 0) / total
    challenge_signals = counts.get("Delayed", 0) + counts.get("Confused", 0) + counts.get("Overwhelmed", 0)

    planning_strength = "High" if st.session_state.study_plan else "Developing"
    if completed_ratio >= 0.6:
        monitoring_awareness = "High"
    elif challenge_signals >= 1 or counts.get("Completed", 0) >= 1:
        monitoring_awareness = "Medium"
    else:
        monitoring_awareness = "Developing"

    if st.session_state.adjusted_plan or challenge_signals >= 2:
        strategy_flexibility = "Developing"
    elif completed_ratio >= 0.7:
        strategy_flexibility = "Medium"
    else:
        strategy_flexibility = "Developing"

    return {
        "Planning strength": planning_strength,
        "Monitoring awareness": monitoring_awareness,
        "Strategy flexibility": strategy_flexibility,
    }


def reflection_report_with_genai(counts: Dict[str, int]) -> Dict[str, object]:
    fallback_profile = learner_profile(counts)
    task = st.session_state.plan_context.get("task_description", "your learning task")
    fallback_summary = (
        f"For {task}, your current SRL profile suggests that you can use planning outputs, "
        "monitoring signals, and strategy adjustments as a repeatable learning cycle."
    )
    fallback_prompts = [
        "今天是什么因素促进或拖慢了你的学习进度？",
        "哪一种策略最有效？为什么？",
        "下一轮学习你会改变什么？",
        "学习任务中的哪一部分变得更清楚了？",
        "什么证据表明你的理解有所提升？",
    ] if current_language() == "中文" else [
        "What helped or slowed your progress today?",
        "Which strategy worked best, and why?",
        "What will you change in the next learning cycle?",
        "Which part of the learning task became clearer?",
        "What evidence shows that your understanding improved?",
    ]
    fallback_insights = (
        [
            "你的反思可以进一步区分努力不足、策略不合适和任务难度较高这三类原因。",
            "下一轮学习应优先保留已经产生可见证据的策略。",
        ]
        if current_language() == "中文"
        else [
            "Your reflection can further separate lack of effort, unsuitable strategy, and task difficulty.",
            "The next cycle should preserve strategies that produced visible evidence of progress.",
        ]
    )
    fallback_next_cycle_plan = (
        [
            "选择一个最小学习目标。",
            "先做主动回忆或自测，再查看材料。",
            "学习结束后记录一个证据和一个仍需澄清的问题。",
        ]
        if current_language() == "中文"
        else [
            "Choose one minimum learning goal.",
            "Use active recall or self-testing before checking materials.",
            "After studying, record one evidence item and one question that still needs clarification.",
        ]
    )
    system_prompt = (
        "You are PlanWise, an LLM-supported SRL reflection coach. Generate learner "
        "profile summaries and reflection prompts for university learning tasks. "
        f"Write all user-facing JSON values in {output_language_name()}. Return valid JSON only."
    )
    user_prompt = f"""
Generate an SRL Reflection-phase report.

Plan context:
{json.dumps(st.session_state.plan_context, ensure_ascii=False)}

Study plan:
{json.dumps(st.session_state.study_plan, ensure_ascii=False)}

Current workflow step:
{json.dumps(current_step(), ensure_ascii=False)}

Monitoring statuses:
{json.dumps(monitoring_status_through_current(), ensure_ascii=False)}

Monitoring details:
{json.dumps(monitoring_details_through_current(), ensure_ascii=False)}

Learner reflection inputs:
{json.dumps(st.session_state.reflection_inputs, ensure_ascii=False)}

Adjusted plan:
{json.dumps(st.session_state.adjusted_plan, ensure_ascii=False)}

Status counts:
{json.dumps(counts, ensure_ascii=False)}

Return JSON with this exact structure:
{{
  "profile": {{
    "Planning strength": "High/Medium/Developing",
    "Monitoring awareness": "High/Medium/Developing",
    "Strategy flexibility": "High/Medium/Developing"
  }},
  "summary": "concise SRL learner profile summary",
  "insights": [
    "reflection insight"
  ],
  "next_cycle_plan": [
    "next-cycle action"
  ],
  "prompts": [
    "reflection prompt"
  ]
}}

Requirements:
- Write the summary and prompts in {output_language_name()}.
- Make the profile consistent with the monitoring evidence.
- Use learner reflection inputs to distinguish attribution, strategy effectiveness, emotion, and next-cycle intent.
- Provide 2 to 4 insights and 3 next_cycle_plan actions.
- Prompts should encourage causal attribution, strategy evaluation, and next-cycle planning.
- Give exactly 5 prompts.
- Do not include Markdown fences.
"""
    try:
        data = call_genai_json(system_prompt, user_prompt)
        profile = data.get("profile", {})
        summary = str(data.get("summary", "")).strip()
        insights = data.get("insights", [])
        next_cycle_plan = data.get("next_cycle_plan", [])
        prompts = data.get("prompts", [])
        if not isinstance(profile, dict) or not profile:
            raise ValueError("GenAI response did not include a profile.")
        if not summary:
            summary = fallback_summary
        if not isinstance(prompts, list) or not prompts:
            prompts = fallback_prompts
        if not isinstance(insights, list) or not insights:
            insights = fallback_insights
        if not isinstance(next_cycle_plan, list) or not next_cycle_plan:
            next_cycle_plan = fallback_next_cycle_plan
        for key, value in fallback_profile.items():
            profile[key] = str(profile.get(key) or value)
        return {
            "profile": profile,
            "summary": summary,
            "insights": [str(item) for item in insights[:4]],
            "next_cycle_plan": [str(item) for item in next_cycle_plan[:3]],
            "prompts": [str(item) for item in prompts[:5]],
            "source": "GenAI-generated via OpenAI API",
        }
    except Exception as exc:
        return {
            "profile": fallback_profile,
            "summary": fallback_summary,
            "insights": fallback_insights,
            "next_cycle_plan": fallback_next_cycle_plan,
            "prompts": fallback_prompts + [f"Fallback note: GenAI reflection report was unavailable ({exc})."],
            "source": "Rule-based fallback after GenAI was unavailable",
        }


def reflection_page() -> None:
    render_header("Reflection")
    st.markdown(f'<div class="phase-label">{t("phase_prefix")}: {phase_label("Reflection")}</div>', unsafe_allow_html=True)
    st.subheader(t("reflection_title"))

    if not st.session_state.study_plan:
        st.info(t("reflection_need_plan"))
        return

    render_workflow_progress()
    active_step = current_step()
    if active_step:
        st.subheader(f"{t('current_step')}: {active_step['day']}")
        render_step_summary(active_step, current_step_index())

    counts = get_status_counts(through_current=True)
    with st.expander(t("reflection_inputs"), expanded=True):
        st.caption(t("reflection_inputs_caption"))
        ref_col1, ref_col2 = st.columns(2)
        existing_inputs = st.session_state.reflection_inputs
        with ref_col1:
            what_worked = st.text_area(
                t("what_worked"),
                value=str(existing_inputs.get("what_worked", "")),
                placeholder=t("what_worked_placeholder"),
                height=90,
            )
            what_did_not_work = st.text_area(
                t("what_did_not_work"),
                value=str(existing_inputs.get("what_did_not_work", "")),
                placeholder=t("what_did_not_work_placeholder"),
                height=90,
            )
            strategy_to_keep = st.text_input(
                t("strategy_to_keep"),
                value=str(existing_inputs.get("strategy_to_keep", "")),
            )
        with ref_col2:
            progress_reason = st.text_area(
                t("progress_reason"),
                value=str(existing_inputs.get("progress_reason", "")),
                placeholder=t("progress_reason_placeholder"),
                height=90,
            )
            emotion_options = TRANSLATIONS["emotion_energy_options"][current_language()]
            existing_emotion = existing_inputs.get("emotion_energy", emotion_options[1])
            emotion_index = emotion_options.index(existing_emotion) if existing_emotion in emotion_options else 1
            emotion_energy = st.selectbox(
                t("emotion_energy"),
                emotion_options,
                index=emotion_index,
            )
            strategy_to_change = st.text_input(
                t("strategy_to_change"),
                value=str(existing_inputs.get("strategy_to_change", "")),
            )
        next_cycle_goal = st.text_area(
            t("next_cycle_goal"),
            value=str(existing_inputs.get("next_cycle_goal", "")),
            placeholder=t("next_cycle_goal_placeholder"),
            height=80,
        )

    st.session_state.reflection_inputs = {
        "what_worked": what_worked.strip(),
        "what_did_not_work": what_did_not_work.strip(),
        "progress_reason": progress_reason.strip(),
        "emotion_energy": emotion_energy,
        "strategy_to_keep": strategy_to_keep.strip(),
        "strategy_to_change": strategy_to_change.strip(),
        "next_cycle_goal": next_cycle_goal.strip(),
    }

    if st.button(t("generate_reflection_report"), type="primary"):
        with st.spinner(t("reflection_spinner")):
            st.session_state.reflection_report = reflection_report_with_genai(counts)

    report = st.session_state.reflection_report
    if report:
        profile = report.get("profile", learner_profile(counts))
        summary = report.get("summary", "")
        insights = report.get("insights", [])
        next_cycle_plan = report.get("next_cycle_plan", [])
        prompts = report.get("prompts", [])
        source = report.get("source", "")
    else:
        profile = learner_profile(counts)
        task = st.session_state.plan_context.get("task_description", "your learning task")
        summary = (
            f"For {task}, your current SRL profile suggests that you can use planning outputs, "
            "monitoring signals, and strategy adjustments as a repeatable learning cycle."
        )
        prompts = [
            "What slowed your progress today?",
            "Which strategy worked best?",
            "What will you change tomorrow?",
            "Which part of the learning task became clearer?",
            "What evidence shows that your understanding improved?",
        ] if current_language() == "English" else [
            "今天是什么因素拖慢了你的学习进度？",
            "哪一种策略最有效？",
            "明天你会改变什么？",
            "学习任务中的哪一部分变得更清楚了？",
            "什么证据表明你的理解有所提升？",
        ]
        insights = (
            [
                "你的反思可以进一步区分努力、策略和任务难度之间的关系。",
                "优先保留能产生可见学习证据的策略。",
            ]
            if current_language() == "中文"
            else [
                "Your reflection can further separate effort, strategy, and task difficulty.",
                "Prioritise keeping strategies that produce visible evidence of learning.",
            ]
        )
        next_cycle_plan = (
            [
                "设定一个最小可完成目标。",
                "先做主动回忆或自测，再查看材料。",
                "学习后记录一个证据和一个仍需澄清的问题。",
            ]
            if current_language() == "中文"
            else [
                "Set one minimum viable learning goal.",
                "Use active recall or self-testing before checking materials.",
                "Record one evidence item and one remaining question after studying.",
            ]
        )
        source = "Rule-based preview before GenAI generation"

    source_caption(source)
    col1, col2, col3 = st.columns(3)
    col1.metric(t("planning_strength"), level_label(profile.get("Planning strength", "Developing")))
    col2.metric(t("monitoring_awareness"), level_label(profile.get("Monitoring awareness", "Developing")))
    col3.metric(t("strategy_flexibility"), level_label(profile.get("Strategy flexibility", "Developing")))

    st.subheader(t("learner_profile_summary"))
    st.write(summary)

    st.subheader(t("reflection_insights"))
    for insight in insights:
        st.markdown(f'<div class="feedback-box">{escape(str(insight))}</div>', unsafe_allow_html=True)

    st.subheader(t("next_cycle_plan"))
    for action in next_cycle_plan:
        st.markdown(f'<div class="warning-box">{escape(str(action))}</div>', unsafe_allow_html=True)

    st.subheader(t("reflection_prompts"))
    for prompt in prompts:
        st.markdown(f"- {prompt}")

    reflection_notes = st.text_area(
        t("reflection_notes"),
        value=str(st.session_state.reflection_inputs.get("reflection_notes", "")),
        placeholder=t("reflection_notes_placeholder"),
        height=140,
    )
    st.session_state.reflection_inputs["reflection_notes"] = reflection_notes.strip()

    can_advance = bool(reflection_notes.strip() or st.session_state.reflection_report)
    total_steps = len(st.session_state.study_plan)
    current_index = current_step_index()
    current_already_completed = current_index in st.session_state.workflow_completed_indices
    if current_index in st.session_state.workflow_completed_indices and current_index >= total_steps - 1:
        st.success(t("all_steps_completed"))
    elif current_already_completed:
        st.info(t("step_already_completed"))
        if st.button(t("back_to_monitoring"), use_container_width=True):
            navigate_to_page("Monitoring")
    elif st.button(t("advance_next_step"), type="primary", disabled=not can_advance):
        mark_current_step_complete()
        set_sidebar_flash("current_step_completed")
        navigate_to_page("Monitoring")
    elif not can_advance:
        st.caption(t("finish_current_step_first"))
    save_current_task()


def has_openai_key() -> bool:
    return bool(get_openai_api_key()) and OpenAI is not None


def get_openai_api_key() -> str:
    try:
        secret_key = st.secrets.get("OPENAI_API_KEY", "")
    except Exception:
        secret_key = ""
    return secret_key or os.getenv("OPENAI_API_KEY", "")


def get_openai_model() -> str:
    try:
        secret_model = st.secrets.get("OPENAI_MODEL", "")
    except Exception:
        secret_model = ""
    return secret_model or os.getenv("OPENAI_MODEL", "gpt-4o-mini")


def get_openai_base_url() -> str:
    try:
        secret_base_url = st.secrets.get("OPENAI_BASE_URL", "")
    except Exception:
        secret_base_url = ""
    return secret_base_url or os.getenv("OPENAI_BASE_URL", "")


def rule_based_coach(user_message: str, context: Dict[str, object]) -> str:
    message = user_message.lower()
    plan_context = context.get("plan_context", context)
    task = plan_context.get(
        "task_description",
        "你的学习任务" if current_language() == "中文" else "your learning task",
    )
    challenge = plan_context.get(
        "challenge",
        "当前学习困难" if current_language() == "中文" else "your current challenge",
    )

    if current_language() == "中文":
        if "读" in user_message or "材料" in user_message or "章节" in user_message:
            return (
                "计划：先明确学习问题，并把材料分成核心、辅助和可选三类。\n\n"
                "监控：学习 30 分钟后，检查自己能否说明每份材料和目标的关系。\n\n"
                "调节：如果材料太多，停止逐份总结，改用综合表或优先级图。\n\n"
                f"反思：回顾哪些材料真正推动了“{task}”。"
            )
        if "结构" in user_message or "整理" in user_message or "组织" in user_message:
            return (
                "计划：先写下一个一句话目标和三个具体里程碑。\n\n"
                "监控：检查每个里程碑是否有行动、进展证据和明确产出。\n\n"
                "调节：如果结构不清楚，先用概念图连接想法、资源和产出。\n\n"
                "反思：识别哪一步让你的理解最明显地增强。"
            )
        if "时间" in user_message or "截止" in user_message or "来不及" in user_message:
            return (
                "计划：定义最小可完成成果，也就是仍能满足目标的最小完整版本。\n\n"
                "监控：区分哪些部分已经可用，哪些只是半成品或笔记。\n\n"
                "调节：压缩可选材料，把时间留给最高价值产出。\n\n"
                "反思：记录哪个时间估计不准确，帮助下一轮计划更现实。"
            )
        return (
            f"计划：把下一步行动连接到“{task}”，并定义一个可见产出。\n\n"
            f"监控：观察“{challenge}”背后的信号，例如延迟、困惑或负荷过高。\n\n"
            "调节：不要只增加学习时间，先调整任务大小、笔记方式或学习顺序。\n\n"
            "反思：写一句你学到了什么，再写一句下一轮要改变什么。"
        )

    if "read" in message or "source" in message or "material" in message or "chapter" in message:
        return (
            "Planning: define the learning question and sort materials into priority groups.\n\n"
            "Monitoring: after 30 minutes, check whether you can explain each material's role in the goal.\n\n"
            "Control: if the material list feels too large, stop summarising everything and build a synthesis or priority map.\n\n"
            f"Reflection: ask which materials actually moved {task} forward."
        )
    if "structure" in message or "organize" in message or "organise" in message or "outline" in message:
        return (
            "Planning: write a one-sentence goal and three concrete milestones before starting deep work.\n\n"
            "Monitoring: check whether each milestone has an action, evidence of progress, and a clear output.\n\n"
            "Control: if the structure feels unclear, create a concept map linking ideas, resources, and outputs.\n\n"
            "Reflection: identify the step where your understanding became strongest."
        )
    if "time" in message or "deadline" in message or "running out" in message:
        return (
            "Planning: define the minimum viable outcome: the smallest complete version that still meets the goal.\n\n"
            "Monitoring: track which parts are ready, which are partial, and which are only notes.\n\n"
            "Control: compress optional materials and protect time for the highest-impact output.\n\n"
            "Reflection: note which time estimate was inaccurate so tomorrow's plan is more realistic."
        )
    if "stuck" in message or "start" in message or "anxiety" in message:
        return (
            "Planning: lower the starting cost by creating a rough first attempt, even if it is incomplete.\n\n"
            "Monitoring: notice whether you are avoiding the task because of uncertainty or perfectionism.\n\n"
            "Control: use a 20-minute focused sprint without editing or over-checking.\n\n"
            "Reflection: record what became easier once you had a visible starting point."
        )

    return (
        f"Planning: connect your next action to {task} and define one visible output.\n\n"
        f"Monitoring: watch for the signal behind {challenge}, such as delay, confusion, or overload.\n\n"
        "Control: adjust the strategy before adding more time. Change the task size, note-taking method, or work sequence.\n\n"
        "Reflection: write one sentence about what you learned and one sentence about what you will change."
    )


def build_coach_context(use_planner_context: bool) -> Dict[str, object]:
    if not use_planner_context:
        return {}
    return {
        "plan_context": st.session_state.plan_context,
        "learning_materials_brief": materials_brief(),
        "study_plan": st.session_state.study_plan,
        "monitoring_status": st.session_state.monitoring_status,
        "monitoring_details": st.session_state.monitoring_details,
        "adjusted_plan": st.session_state.adjusted_plan,
        "control_context": st.session_state.control_context,
        "reflection_inputs": st.session_state.reflection_inputs,
    }


def openai_coach(
    user_message: str,
    context: Dict[str, object],
    coach_settings: Dict[str, object],
) -> Tuple[str, str]:
    if not has_openai_key():
        return rule_based_coach(user_message, context), "Rule-based fallback because API key is unavailable"

    try:
        reply = call_genai_text(
            (
                "You are PlanWise, a metacognitive learning coach for university students. "
                "Give concise study and learning support aligned with Zimmerman's SRL cycle: "
                f"Planning, Monitoring, Control, and Reflection. Write in {output_language_name()}. "
                "Do not act like a generic chatbot. Use the learner's settings to choose the scaffolding style."
            ),
            (
                f"Coach settings: {json.dumps(coach_settings, ensure_ascii=False)}\n\n"
                f"Current learner context: {json.dumps(context, ensure_ascii=False)}\n\n"
                f"Student message: {user_message}\n\n"
                f"""
Return a structured coaching response in {output_language_name()} with:
1. A short diagnosis of the learner's current barrier.
2. Guidance aligned with the selected SRL focus.
3. Two or three immediate next actions.
4. One metacognitive check question.

If the response style is Socratic questions, use more questions and fewer direct instructions.
If urgency is high, prioritise short-term rescue actions.
"""
            ),
            temperature=0.4,
        )
        if not reply:
            raise ValueError("GenAI returned an empty coach response.")
        return reply, "GenAI-generated via OpenAI API"
    except Exception as exc:
        return (
            rule_based_coach(user_message, context)
            + f"\n\nNote: OpenAI response was unavailable, so PlanWise used the rule-based coach. Technical detail: {exc}",
            "Rule-based fallback after GenAI was unavailable",
        )


def coach_page() -> None:
    render_header("AI Coach")
    support_panel = "SRL Support Panel" if current_language() == "English" else "自我调节学习支持面板"
    st.markdown(f'<div class="phase-label">{support_panel}</div>', unsafe_allow_html=True)
    st.subheader(t("coach_panel"))
    st.write(t("coach_intro"))

    status = t("api_detected") if has_openai_key() else t("fallback_active")
    st.caption(status)
    render_material_upload_panel("coach")

    with st.expander(t("coach_settings"), expanded=True):
        st.caption(t("coach_settings_caption"))
        setting_col1, setting_col2 = st.columns(2)
        situation_options = TRANSLATIONS["learning_situation_options"][current_language()]
        focus_options = TRANSLATIONS["srl_focus_options"][current_language()]
        style_options = TRANSLATIONS["coach_response_style_options"][current_language()]
        urgency_options = TRANSLATIONS["urgency_options"][current_language()]
        with setting_col1:
            learning_situation = st.selectbox(t("learning_situation"), situation_options)
            st.markdown(f"**{t('srl_focus')}**")
            focus_cols = st.columns(2)
            selected_focus = []
            for index, option in enumerate(focus_options):
                checked = focus_cols[index % 2].checkbox(
                    option,
                    value=index < 2,
                    key=f"coach_srl_focus_{index}",
                )
                if checked:
                    selected_focus.append(option)
            srl_focus = selected_focus or [focus_options[0]]
        with setting_col2:
            response_style = st.selectbox(t("coach_response_style"), style_options)
            urgency = st.selectbox(t("urgency"), urgency_options)
        use_planner_context = st.checkbox(t("use_planner_context"), value=True)
        st.caption(t("coach_context_note"))

    coach_settings = {
        "learning_situation": learning_situation,
        "srl_focus": srl_focus,
        "response_style": response_style,
        "urgency": urgency,
        "use_planner_context": use_planner_context,
    }

    example_prompt = st.selectbox(
        t("example_prompts"),
        [
            "",
            t("example_prompt_1"),
            t("example_prompt_2"),
            t("example_prompt_3"),
        ],
    )
    user_message = st.text_area(
        t("student_concern"),
        value=example_prompt,
        placeholder=t("student_concern_placeholder"),
        height=110,
    )

    if st.button(t("ask_coach"), type="primary"):
        if not user_message.strip():
            st.warning(t("coach_warning"))
        else:
            with st.spinner(t("coach_spinner")):
                learner_context = build_coach_context(use_planner_context)
                reply, source = openai_coach(
                    user_message.strip(),
                    learner_context,
                    coach_settings,
                )
            st.session_state.coach_history.append(
                {
                    "user": user_message.strip(),
                    "assistant": reply,
                    "source": source,
                    "settings": coach_settings,
                }
            )
            save_current_task(create_if_missing=True)

    if st.session_state.coach_history:
        if st.button(t("clear_coach_history")):
            st.session_state.coach_history = []
            save_current_task()
            st.success(t("coach_history_cleared"))
            st.rerun()
        st.subheader(t("coaching_history"))
        for turn in reversed(st.session_state.coach_history[-5:]):
            safe_user = escape(turn["user"])
            safe_assistant = escape(turn["assistant"]).replace("\n", "<br>")
            safe_source = escape(translated_source(turn.get("source", "")))
            settings = turn.get("settings", {})
            safe_settings = escape(
                " | ".join(
                    str(value)
                    for key, value in settings.items()
                    if key != "use_planner_context" and value
                )
            )
            st.markdown(
                f"""
                <div class="coach-box">
                    <small>{safe_source}</small><br>
                    <small>{safe_settings}</small><br>
                    <strong>{t("student")}:</strong> {safe_user}<br><br>
                    <strong>{t("planwise_coach")}:</strong><br>{safe_assistant}
                </div>
                """,
                unsafe_allow_html=True,
            )


def sidebar_navigation() -> str:
    st.sidebar.title(t("nav_title"))
    st.sidebar.caption(t("nav_caption"))
    st.session_state.language = st.sidebar.radio(
        t("language"),
        LANGUAGE_OPTIONS,
        index=LANGUAGE_OPTIONS.index(st.session_state.get("language", "English")),
        horizontal=True,
    )
    render_sidebar_flash()

    st.sidebar.divider()
    st.sidebar.subheader(t("task_memory"))
    st.sidebar.caption(f"{t('saved_task_count')}: {len(st.session_state.tasks)}")
    if st.sidebar.button(t("open_task_library"), use_container_width=True):
        navigate_to_page("Task Library")
    if st.session_state.tasks:
        current_task = st.session_state.tasks.get(st.session_state.current_task_id, {})
        if current_task:
            st.sidebar.caption(f"{t('current_task')}: {current_task.get('title', t('unsaved_task'))}")
    else:
        st.sidebar.caption(t("no_saved_tasks"))

    task_col1, task_col2 = st.sidebar.columns(2)
    if task_col1.button(t("new_task"), use_container_width=True):
        reset_learning_state()
        set_sidebar_flash("new_task_ready")
        st.rerun()
    if task_col2.button(t("save_task"), use_container_width=True):
        current_task_name = ""
        current_task_id = st.session_state.get("current_task_id", "")
        if current_task_id:
            current_task_name = str(st.session_state.get(f"rename_title_{current_task_id}", "")).strip()
        saved_id = save_current_task(create_if_missing=True, custom_title=current_task_name)
        if saved_id:
            sync_task_selection(saved_id)
            set_sidebar_flash("task_saved")
            st.rerun()
        else:
            set_sidebar_flash("nothing_to_save", level="warning")
            st.rerun()

    if st.session_state.current_task_id:
        current_task = st.session_state.tasks.get(st.session_state.current_task_id, {})
        updated_at = str(current_task.get("updated_at", ""))[:16].replace("T", " ")
        if updated_at:
            st.sidebar.caption(f"{t('autosaved')}: {updated_at}")

    st.sidebar.divider()
    pages = [
        "Dashboard",
        "Task Library",
        "Task Details",
        "Planning",
        "Today Learning",
        "Monitoring",
        "Control",
        "Reflection",
        "AI Coach",
        "Reports",
        "How to Use",
    ]
    pending_page = st.session_state.get("pending_page_selection")
    if pending_page in pages:
        st.session_state.page_selector = pending_page
        st.session_state.pending_page_selection = None
    page = st.sidebar.radio(
        t("go_to"),
        pages,
        format_func=phase_label,
        key="page_selector",
    )

    if st.session_state.study_plan:
        st.sidebar.divider()
        st.sidebar.subheader(t("current_task"))
        st.sidebar.write(st.session_state.plan_context.get("task_description", "Learning task"))
        st.sidebar.caption(f"{t('deadline')}: {st.session_state.plan_context.get('deadline', 'Not set')}")
    return page


def main() -> None:
    st.set_page_config(
        page_title="PlanWise SRL Planner",
        page_icon="PW",
        layout="wide",
        initial_sidebar_state="expanded",
    )
    initialize_state()
    inject_styles()
    page = sidebar_navigation()

    if page == "Dashboard":
        dashboard_page()
    elif page == "Task Library":
        task_library_page()
    elif page == "Task Details":
        task_details_page()
    elif page == "Planning":
        planning_page()
    elif page == "Today Learning":
        today_learning_page()
    elif page == "Monitoring":
        monitoring_page()
    elif page == "Control":
        control_page()
    elif page == "Reflection":
        reflection_page()
    elif page == "AI Coach":
        coach_page()
    elif page == "Reports":
        reports_page()
    else:
        guide_page()


if __name__ == "__main__":
    main()
